import os, datetime
from fastapi import FastAPI, Request, Form, Depends, HTTPException
from fastapi.responses import RedirectResponse, HTMLResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from starlette.middleware.sessions import SessionMiddleware

from . import db
from .auth import (AUTH_MODE, domain_ok, current_user, require_user, require_editor,
                   require_admin, get_role, ALLOWED_DOMAIN)
from . import auth as _auth

app = FastAPI(title="Qcells AVL Manager")
# https_only stops the session cookie ever being sent in clear; set SECURE_COOKIES=1
# once the app is behind TLS. Left off by default so plain-http local dev works.
app.add_middleware(SessionMiddleware,
                   secret_key=os.environ.get("SECRET_KEY", "change-me-in-prod"),
                   https_only=os.environ.get("SECURE_COOKIES", "") == "1",
                   same_site="lax")
BASE = os.path.dirname(__file__)
app.mount("/static", StaticFiles(directory=os.path.join(BASE, "static")), name="static")
templates = Jinja2Templates(directory=os.path.join(BASE, "templates"))

def _css_v():
    """Stylesheet mtime, so a CSS edit is not hidden behind a browser cache."""
    try:
        return int(os.path.getmtime(os.path.join(BASE, "static", "style.css")))
    except OSError:
        return 0

templates.env.globals["css_v"] = _css_v

db.init_db()

for _w in _auth.startup_warnings():
    print(f"WARNING: {_w}", flush=True)

@app.get("/healthz")
def healthz():
    """Unauthenticated liveness probe for the host's health check."""
    return {"ok": True}

def now():
    return datetime.datetime.now().isoformat(timespec="seconds")

# ---------------- auth routes ----------------
@app.get("/login", response_class=HTMLResponse)
def login_page(request: Request):
    return templates.TemplateResponse(request, "login.html", {"user": None, "mode": AUTH_MODE,
        "domain": ALLOWED_DOMAIN, "warnings": _auth.startup_warnings(),
        "needs_password": AUTH_MODE == "shared"})

def _client_ip(request: Request):
    fwd = request.headers.get("X-Forwarded-For", "")
    return (fwd.split(",")[0].strip() if fwd else None) or \
           (request.client.host if request.client else "unknown")

@app.post("/login/dev")
def login_dev(request: Request, name: str = Form(...), email: str = Form(...),
              password: str = Form("")):
    """Form login for dev and shared modes; shared additionally needs the password."""
    ip = _client_ip(request)
    if not _auth.form_login_enabled():
        return RedirectResponse("/login", status_code=303)
    if _auth.login_blocked(ip):
        return RedirectResponse("/login?error=locked", status_code=303)
    if not domain_ok(email):
        _auth.note_login_failure(ip)
        return RedirectResponse("/login?error=domain", status_code=303)
    if AUTH_MODE == "shared":
        if not _auth.shared_password_configured():
            return RedirectResponse("/login?error=unconfigured", status_code=303)
        if not _auth.shared_password_ok(password):
            _auth.note_login_failure(ip)
            db.log(email.lower().strip(), "login:failed", f"bad shared password from {ip}")
            return RedirectResponse("/login?error=password", status_code=303)
    if not _auth.user_known(email):
        _auth.note_login_failure(ip)
        db.log(email.lower().strip(), "login:refused", f"not on the user list, from {ip}")
        return RedirectResponse("/login?error=notinvited", status_code=303)
    _auth.clear_login_failures(ip)
    request.session["user"] = {"email": email.lower().strip(), "name": name.strip()}
    _touch_user(request.session["user"])
    return RedirectResponse("/", status_code=303)

@app.get("/login/sso")
async def login_sso(request: Request):
    from .auth import oauth_client
    redirect_uri = request.url_for("auth_callback")
    return await oauth_client().authorize_redirect(request, redirect_uri)

@app.get("/auth/callback")
async def auth_callback(request: Request):
    from .auth import oauth_client
    token = await oauth_client().authorize_access_token(request)
    info = token.get("userinfo") or {}
    email = (info.get("email") or info.get("preferred_username") or "").lower()
    if not domain_ok(email):
        return RedirectResponse("/login?error=domain", status_code=303)
    if not _auth.user_known(email):
        db.log(email, "login:refused", "not on the user list (sso)")
        return RedirectResponse("/login?error=notinvited", status_code=303)
    request.session["user"] = {"email": email, "name": info.get("name", email.split("@")[0])}
    _touch_user(request.session["user"])
    return RedirectResponse("/", status_code=303)

@app.get("/logout")
def logout(request: Request):
    request.session.clear()
    return RedirectResponse("/login", status_code=303)

def _touch_user(user):
    # A login is a good moment to check whether today's snapshot exists: it is
    # infrequent, it happens before anyone edits anything, and it needs no
    # scheduler. auto_backup() is a no-op if one was taken recently.
    db.auto_backup()
    c = db.conn()
    if c.execute("SELECT COUNT(*) FROM users").fetchone()[0] == 0:
        c.execute("INSERT INTO users(email, name, role, last_login) VALUES(?,?, 'admin', ?)",
                  (user["email"], user["name"], now()))
        c.commit(); c.close(); return
    c.execute("INSERT INTO users(email, name, last_login) VALUES(?,?,?) "
              "ON CONFLICT(email) DO UPDATE SET name=excluded.name, last_login=excluded.last_login",
              (user["email"], user["name"], now()))
    c.commit(); c.close()

# ---------------- dashboard ----------------
@app.get("/", response_class=HTMLResponse)
def dashboard(request: Request, show_eol: int = 0, user=Depends(require_user)):
    c = db.conn()
    avls = c.execute("SELECT * FROM avls WHERE active=1 ORDER BY id").fetchall()
    q = "SELECT * FROM products WHERE active=1 "
    if not show_eol:
        q += "AND lifecycle != 'EOL' "
    products = c.execute(q + "ORDER BY category, id").fetchall()
    grid = {}
    for row in c.execute("SELECT product_id, avl_id, status, note, updated_by, updated_at FROM listings"):
        grid[(row["product_id"], row["avl_id"])] = row
    am_map = {}
    for row in c.execute("SELECT a.avl_id, p.name FROM assignments a JOIN people p ON p.id=a.person_id "
                         "WHERE a.role LIKE 'Account Manager%' AND a.ended_at IS NULL AND a.avl_id IS NOT NULL"):
        am_map.setdefault(row["avl_id"], []).append(row["name"])
    avls = [dict(a) | {"account_manager": ", ".join(am_map.get(a["id"], [])) or a["account_manager"]} for a in avls]
    avls = [type("A", (), a) for a in avls]
    c.close()
    return templates.TemplateResponse(request, "dashboard.html", {
        "user": user, "avls": avls, "products": products,
        "grid": grid, "statuses": db.STATUSES, "show_eol": show_eol})

@app.post("/listing")
def update_listing(request: Request, product_id: int = Form(...), avl_id: int = Form(...),
                   status: str = Form(...), user=Depends(require_editor)):
    if status not in db.STATUSES:
        return RedirectResponse("/", status_code=303)
    c = db.conn()
    prev = c.execute("SELECT status FROM listings WHERE product_id=? AND avl_id=?",
                     (product_id, avl_id)).fetchone()
    old_status = prev["status"] if prev else None
    c.execute("INSERT INTO listings(product_id, avl_id, status, updated_by, updated_at) VALUES(?,?,?,?,?) "
              "ON CONFLICT(product_id, avl_id) DO UPDATE SET status=excluded.status, "
              "updated_by=excluded.updated_by, updated_at=excluded.updated_at",
              (product_id, avl_id, status, user["email"], now()))
    if old_status != status:
        c.execute("INSERT INTO status_history(product_id, avl_id, old_status, new_status, changed_by, ts) "
                  "VALUES(?,?,?,?,?,?)", (product_id, avl_id, old_status, status, user["email"], now()))
    pname = c.execute("SELECT name FROM products WHERE id=?", (product_id,)).fetchone()["name"]
    aname = c.execute("SELECT name FROM avls WHERE id=?", (avl_id,)).fetchone()["name"]
    c.commit(); c.close()
    db.log(user["email"], "status", f"{pname} @ {aname} -> {status}",
           avl_id=avl_id, product_id=product_id, entity="listing")
    return RedirectResponse("/", status_code=303)

# ---------------- manage products / AVLs / managers ----------------
@app.get("/manage", response_class=HTMLResponse)
def manage(request: Request, user=Depends(require_user)):
    c = db.conn()
    avls = c.execute("SELECT * FROM avls ORDER BY active DESC, id").fetchall()
    products = c.execute("SELECT * FROM products ORDER BY active DESC, category, id").fetchall()
    people = c.execute("SELECT id, name, org FROM people WHERE active=1 ORDER BY name").fetchall()
    AM, SCR, PT = db.ROLES
    # Who currently holds each seat, so the pickers come up pre-selected.
    holders = {
        "avl": {a["id"]: {"am": [r["person_id"] for r in db.role_holders(c, AM, avl_id=a["id"])],
                          "scr": [r["person_id"] for r in db.role_holders(c, SCR, avl_id=a["id"])]}
                for a in avls},
        "product": {p["id"]: [r["person_id"] for r in db.role_holders(c, PT, product_id=p["id"])]
                    for p in products},
    }
    AM_, SCR_, PT_ = db.ROLES
    # Options are built per row so somebody already holding a seat stays listed
    # even if their team would not normally be offered for it.
    by_id = {p["id"]: p for p in people}
    # Reps render as chips, so the row needs the assigned people in order and the
    # ones still available to add.
    prod_reps = {p["id"]: [by_id[i] for i in holders["product"].get(p["id"], []) if i in by_id]
                 for p in products}
    prod_avail = {p["id"]: [(lbl, [r for r in folk
                                   if r["id"] not in holders["product"].get(p["id"], [])])
                            for lbl, folk in db.role_options(people, PT_)]
                  for p in products}
    prod_avail = {k: [(lbl, folk) for lbl, folk in v if folk] for k, v in prod_avail.items()}
    avl_opts = {a["id"]: {
        "am": db.role_options(people, AM_, holders["avl"].get(a["id"], {}).get("am", [])),
        "scr": db.role_options(people, SCR_, holders["avl"].get(a["id"], {}).get("scr", [])),
    } for a in avls}
    new_prod_opts = db.role_options(people, PT_)
    new_avl_opts = db.role_options(people, AM_)
    c.close()
    return templates.TemplateResponse(request, "manage.html", {"user": user,
                            "avls": avls, "products": products, "categories": db.CATEGORIES,
                            "people": people, "holders": holders,
                            "prod_reps": prod_reps, "prod_avail": prod_avail,
                            "avl_opts": avl_opts,
                            "new_prod_opts": new_prod_opts, "new_avl_opts": new_avl_opts})

@app.post("/products/add")
def add_product(request: Request, name: str = Form(...), category: str = Form(...),
                person_id: str = Form(""), launch_status: str = Form("Released"),
                lifecycle: str = Form("Active"), user=Depends(require_editor)):
    if lifecycle not in ("Roadmap", "Active", "EOL"):
        lifecycle = "Active"
    _, _, PT = db.ROLES
    c = db.conn()
    c.execute("INSERT OR IGNORE INTO products(name, category, launch_status, lifecycle) "
              "VALUES(?,?,?,?)", (name.strip(), category, launch_status.strip(), lifecycle))
    row = c.execute("SELECT id FROM products WHERE name=?", (name.strip(),)).fetchone()
    if row:
        picked = [int(person_id)] if person_id.strip().isdigit() and int(person_id) > 0 else []
        db.set_role_holders(c, PT, picked, product_id=row["id"], actor=user["email"])
        db.refresh_rep_cache(c, product_id=row["id"])
    c.commit(); c.close()
    db.log(user["email"], "product:add", name)
    return RedirectResponse("/manage", status_code=303)

@app.post("/products/{pid}/toggle")
def toggle_product(pid: int, request: Request, user=Depends(require_editor)):
    c = db.conn()
    c.execute("UPDATE products SET active = 1 - active WHERE id=?", (pid,))
    name = c.execute("SELECT name, active FROM products WHERE id=?", (pid,)).fetchone()
    c.commit(); c.close()
    db.log(user["email"], "product:toggle", f"{name['name']} active={name['active']}")
    return RedirectResponse("/manage", status_code=303)

def _set_reps(c, pid, person_ids, actor):
    _, _, PT = db.ROLES
    res = db.set_role_holders(c, PT, person_ids, product_id=pid, actor=actor)
    db.refresh_rep_cache(c, product_id=pid)
    return res

@app.post("/products/{pid}/reps/add")
def add_rep(pid: int, request: Request, person_id: int = Form(...), user=Depends(require_editor)):
    _, _, PT = db.ROLES
    c = db.conn()
    cur = [r["person_id"] for r in db.role_holders(c, PT, product_id=pid)]
    if person_id not in cur:
        added, _ = _set_reps(c, pid, cur + [person_id], user["email"])
        name = c.execute("SELECT name FROM products WHERE id=?", (pid,)).fetchone()["name"]
        c.commit()
        if added:
            db.log(user["email"], "product:reps", f"{name}: +{', '.join(added)}")
    c.close()
    return RedirectResponse("/manage", status_code=303)

@app.post("/products/{pid}/reps/remove")
def remove_rep(pid: int, request: Request, person_id: int = Form(...), user=Depends(require_editor)):
    _, _, PT = db.ROLES
    c = db.conn()
    cur = [r["person_id"] for r in db.role_holders(c, PT, product_id=pid)]
    if person_id in cur:
        _, removed = _set_reps(c, pid, [i for i in cur if i != person_id], user["email"])
        name = c.execute("SELECT name FROM products WHERE id=?", (pid,)).fetchone()["name"]
        c.commit()
        if removed:
            db.log(user["email"], "product:reps", f"{name}: -{', '.join(removed)}")
    c.close()
    return RedirectResponse("/manage", status_code=303)

@app.post("/products/{pid}/reps")
def set_reps(pid: int, request: Request, person_ids: list[int] = Form(default=[]),
             user=Depends(require_editor)):
    """Product/Technical reps come from the Team roster, never free text."""
    _, _, PT = db.ROLES
    c = db.conn()
    added, removed = db.set_role_holders(c, PT, person_ids, product_id=pid, actor=user["email"])
    db.refresh_rep_cache(c, product_id=pid)
    name = c.execute("SELECT name FROM products WHERE id=?", (pid,)).fetchone()["name"]
    c.commit(); c.close()
    if added or removed:
        db.log(user["email"], "product:reps",
               f"{name}: +{', '.join(added) or '-'} / -{', '.join(removed) or '-'}")
    return RedirectResponse("/manage", status_code=303)

@app.post("/avls/add")
def add_avl(request: Request, name: str = Form(...), account_manager_id: str = Form(""),
            sr_commercial_rep_id: str = Form(""), user=Depends(require_editor)):
    AM, SCR, _ = db.ROLES
    c = db.conn()
    c.execute("INSERT OR IGNORE INTO avls(name) VALUES(?)", (name.strip(),))
    row = c.execute("SELECT id FROM avls WHERE name=?", (name.strip(),)).fetchone()
    if row:
        for role, raw in ((AM, account_manager_id), (SCR, sr_commercial_rep_id)):
            db.set_role_holders(c, role, [int(raw)] if raw.strip().isdigit() else [],
                                avl_id=row["id"], actor=user["email"])
        db.refresh_rep_cache(c, avl_id=row["id"])
    c.commit(); c.close()
    db.log(user["email"], "avl:add", name)
    return RedirectResponse("/manage", status_code=303)

@app.post("/avls/{aid}/toggle")
def toggle_avl(aid: int, request: Request, user=Depends(require_editor)):
    c = db.conn()
    c.execute("UPDATE avls SET active = 1 - active WHERE id=?", (aid,))
    row = c.execute("SELECT name, active FROM avls WHERE id=?", (aid,)).fetchone()
    c.commit(); c.close()
    db.log(user["email"], "avl:toggle", f"{row['name']} active={row['active']}")
    return RedirectResponse("/manage", status_code=303)

@app.post("/avls/{aid}/managers")
def set_managers(aid: int, request: Request, account_manager_id: str = Form(""),
                 sr_commercial_rep_id: str = Form(""), user=Depends(require_editor)):
    """AM and Sr. Commercial Rep come from the Team roster; blank = seat vacant."""
    AM, SCR, _ = db.ROLES
    c = db.conn()
    changes = []
    for role, raw in ((AM, account_manager_id), (SCR, sr_commercial_rep_id)):
        picked = [int(raw)] if raw.strip().isdigit() else []
        added, removed = db.set_role_holders(c, role, picked, avl_id=aid, actor=user["email"])
        if added or removed:
            changes.append(f"{role.split(' (')[0]}: +{', '.join(added) or '-'} / -{', '.join(removed) or '-'}")
    db.refresh_rep_cache(c, avl_id=aid)
    name = c.execute("SELECT name FROM avls WHERE id=?", (aid,)).fetchone()["name"]
    c.commit(); c.close()
    if changes:
        db.log(user["email"], "avl:managers", f"{name} " + "; ".join(changes))
    return RedirectResponse("/manage", status_code=303)

# ---------------- activity: audit log + status changelog in one feed ----------------
def _activity(c, avl=0, product=0, kind="", who="", month="", q="", limit=400):
    """Merged, filtered event feed. Status changes come from status_history so the
    from/to is preserved; everything else from the audit log."""
    rows = []
    sql = ("SELECT a.*, av.name AS avl_name, p.name AS product_name FROM audit a "
           "LEFT JOIN avls av ON av.id=a.avl_id LEFT JOIN products p ON p.id=a.product_id WHERE 1=1 ")
    args = []
    if avl:
        sql += "AND a.avl_id=? "; args.append(avl)
    if product:
        sql += "AND a.product_id=? "; args.append(product)
    if who:
        sql += "AND a.user_email=? "; args.append(who)
    if month:
        sql += "AND substr(a.ts,1,7)=? "; args.append(month)
    if q.strip():
        sql += "AND (a.detail LIKE ? OR a.action LIKE ?) "; args += [f"%{q.strip()}%"] * 2
    for r in c.execute(sql + "ORDER BY a.id DESC LIMIT ?", args + [limit]):
        rows.append({"ts": r["ts"], "who": r["user_email"], "kind": db.activity_type(r["action"]),
                     "action": r["action"], "detail": r["detail"] or "",
                     "avl_id": r["avl_id"], "avl": r["avl_name"],
                     "product_id": r["product_id"], "product": r["product_name"],
                     "from": None, "to": None})

    sql = ("SELECT h.*, av.name AS avl_name, p.name AS product_name FROM status_history h "
           "JOIN avls av ON av.id=h.avl_id JOIN products p ON p.id=h.product_id WHERE 1=1 ")
    args = []
    if avl:
        sql += "AND h.avl_id=? "; args.append(avl)
    if product:
        sql += "AND h.product_id=? "; args.append(product)
    if who:
        sql += "AND h.changed_by=? "; args.append(who)
    if month:
        sql += "AND substr(h.ts,1,7)=? "; args.append(month)
    if q.strip():
        sql += "AND (h.new_status LIKE ? OR h.note LIKE ?) "; args += [f"%{q.strip()}%"] * 2
    for r in c.execute(sql + "ORDER BY h.id DESC LIMIT ?", args + [limit]):
        rows.append({"ts": r["ts"], "who": r["changed_by"], "kind": "status",
                     "action": "status", "detail": r["note"] or "",
                     "avl_id": r["avl_id"], "avl": r["avl_name"],
                     "product_id": r["product_id"], "product": r["product_name"],
                     "from": r["old_status"], "to": r["new_status"]})

    rows.sort(key=lambda r: r["ts"], reverse=True)
    if kind:
        rows = [r for r in rows if r["kind"] == kind]
    return rows[:limit]

@app.get("/activity", response_class=HTMLResponse)
def activity(request: Request, avl: int = 0, product: int = 0, kind: str = "", who: str = "",
             month: str = "", q: str = "", user=Depends(require_user)):
    c = db.conn()
    rows = _activity(c, avl, product, kind, who, month, q)
    avls = c.execute("SELECT id, name FROM avls ORDER BY name").fetchall()
    products = c.execute("SELECT id, name FROM products ORDER BY name").fetchall()
    people = [r[0] for r in c.execute(
        "SELECT DISTINCT user_email FROM audit WHERE COALESCE(user_email,'')<>'' "
        "UNION SELECT DISTINCT changed_by FROM status_history WHERE COALESCE(changed_by,'')<>'' "
        "ORDER BY 1")]
    months = [r[0] for r in c.execute(
        "SELECT DISTINCT substr(ts,1,7) FROM audit "
        "UNION SELECT DISTINCT substr(ts,1,7) FROM status_history ORDER BY 1 DESC")]
    c.close()
    return templates.TemplateResponse(request, "activity.html", {"user": user, "rows": rows,
        "avls": avls, "products": products, "people": people, "months": months,
        "types": db.ACTIVITY_TYPES, "sel": {"avl": avl, "product": product, "kind": kind,
                                            "who": who, "month": month, "q": q}})

@app.get("/activity.csv")
def activity_csv(request: Request, avl: int = 0, product: int = 0, kind: str = "", who: str = "",
                 month: str = "", q: str = "", user=Depends(require_user)):
    c = db.conn()
    rows = _activity(c, avl, product, kind, who, month, q, limit=5000)
    c.close()
    return _csv_response([[r["ts"], r["kind"], r["action"], r["avl"] or "", r["product"] or "",
                           (f"{r['from'] or '-'} -> {r['to']}" if r["to"] else r["detail"]),
                           r["who"]] for r in rows],
                         ["When", "Type", "Event", "TPO AVL", "Product", "Detail", "Who"],
                         "activity.csv")

# The two old pages are now views of the same feed; keep the URLs working.
@app.get("/history")
def history_redirect(request: Request, avl: int = 0, month: str = ""):
    return RedirectResponse(f"/activity?kind=status&avl={avl}&month={month}", status_code=307)

@app.get("/audit")
def audit_redirect(request: Request):
    return RedirectResponse("/activity", status_code=307)

# ---------------- interaction log ----------------
@app.get("/calls", response_class=HTMLResponse)
def calls(request: Request, avl: int = 0, q: str = "", edit: int = 0,
          user=Depends(require_user)):
    c = db.conn()
    avls = c.execute("SELECT * FROM avls WHERE active=1 ORDER BY name").fetchall()
    people = c.execute("SELECT id, name, org FROM people WHERE active=1 ORDER BY name").fetchall()
    sql = "SELECT calls.*, avls.name AS avl_name FROM calls JOIN avls ON avls.id=calls.avl_id WHERE 1=1 "
    params = []
    if avl:
        sql += "AND avl_id=? "
        params.append(avl)
    if q.strip():
        sql += ("AND (topics LIKE ? OR outcomes LIKE ? OR owner_due LIKE ? "
                "OR qcells_attendees LIKE ? OR tpo_attendees LIKE ?) ")
        params += [f"%{q.strip()}%"] * 5
    rows = c.execute(sql + "ORDER BY call_date DESC, calls.id DESC LIMIT 300", params).fetchall()
    # What each logged call already has selected, so the edit form comes up filled in.
    picked = {r["id"]: {"qcells": [], "tpo": [], "other_q": [], "other_t": []} for r in rows}
    if rows:
        qs = ",".join("?" * len(rows))
        for at in c.execute(f"SELECT * FROM call_attendees WHERE call_id IN ({qs})",
                            [r["id"] for r in rows]):
            pk = picked[at["call_id"]]
            if at["person_id"]:
                pk["qcells"].append(at["person_id"])
            elif at["contact_id"]:
                pk["tpo"].append(at["contact_id"])
            else:
                pk["other_q" if at["side"] == "qcells" else "other_t"].append(at["name"])
    # avl -> its active contacts, driving the TPO attendee picker
    by_avl = {}
    for ct in c.execute("SELECT id, avl_id, name, role FROM contacts WHERE active=1 ORDER BY name"):
        by_avl.setdefault(str(ct["avl_id"]), []).append(
            {"id": ct["id"], "name": ct["name"], "role": ct["role"]})
    c.close()
    return templates.TemplateResponse(request, "calls.html", {"user": user, "avls": avls,
        "rows": rows, "sel": avl, "q": q, "people": people, "contacts_by_avl": by_avl,
        "picked": picked, "call_types": db.CALL_TYPES, "edit": edit,
        "today": datetime.date.today().isoformat()})

def _save_attendees(c, cid, qcells_person_ids, tpo_contact_ids, qcells_other, tpo_other):
    db.set_call_attendees(c, cid, "qcells", person_ids=qcells_person_ids, other=qcells_other)
    db.set_call_attendees(c, cid, "tpo", contact_ids=tpo_contact_ids, other=tpo_other)

@app.post("/calls/add")
def add_call(request: Request, avl_id: int = Form(...), call_date: str = Form(...),
             call_type: str = Form("Joint"),
             qcells_person_ids: list[int] = Form(default=[]),
             tpo_contact_ids: list[int] = Form(default=[]),
             qcells_other: str = Form(""), tpo_other: str = Form(""),
             topics: str = Form(""), outcomes: str = Form(""),
             owner_due: str = Form(""), user=Depends(require_editor)):
    c = db.conn()
    c.execute("INSERT INTO calls(avl_id, call_date, call_type, topics, outcomes, owner_due, "
              "created_by, created_at) VALUES(?,?,?,?,?,?,?,?)",
              (avl_id, call_date, call_type, topics.strip(), outcomes.strip(),
               owner_due.strip(), user["email"], now()))
    cid = c.execute("SELECT last_insert_rowid() AS id").fetchone()["id"]
    _save_attendees(c, cid, qcells_person_ids, tpo_contact_ids, qcells_other, tpo_other)
    aname = c.execute("SELECT name FROM avls WHERE id=?", (avl_id,)).fetchone()["name"]
    c.commit(); c.close()
    db.log(user["email"], "call:add", f"{call_date} ({call_type})",
           avl_id=avl_id, entity="call", entity_id=cid)
    return RedirectResponse(f"/calls?avl={avl_id}", status_code=303)

@app.post("/calls/{cid}/save")
def save_call(cid: int, request: Request, avl_id: int = Form(...), call_date: str = Form(...),
              call_type: str = Form("Joint"),
              qcells_person_ids: list[int] = Form(default=[]),
              tpo_contact_ids: list[int] = Form(default=[]),
              qcells_other: str = Form(""), tpo_other: str = Form(""),
              topics: str = Form(""), outcomes: str = Form(""),
              owner_due: str = Form(""), user=Depends(require_editor)):
    c = db.conn()
    if not c.execute("SELECT 1 FROM calls WHERE id=?", (cid,)).fetchone():
        c.close()
        return RedirectResponse("/calls", status_code=303)
    c.execute("UPDATE calls SET avl_id=?, call_date=?, call_type=?, topics=?, outcomes=?, "
              "owner_due=? WHERE id=?", (avl_id, call_date, call_type, topics.strip(),
                                         outcomes.strip(), owner_due.strip(), cid))
    _save_attendees(c, cid, qcells_person_ids, tpo_contact_ids, qcells_other, tpo_other)
    aname = c.execute("SELECT name FROM avls WHERE id=?", (avl_id,)).fetchone()["name"]
    c.commit(); c.close()
    db.log(user["email"], "call:save", f"{call_date}", avl_id=avl_id, entity="call", entity_id=cid)
    return RedirectResponse(f"/calls?avl={avl_id}", status_code=303)

@app.post("/calls/{cid}/delete")
def delete_call(cid: int, request: Request, user=Depends(require_editor)):
    c = db.conn()
    row = c.execute("SELECT calls.*, a.name AS avl_name FROM calls JOIN avls a ON a.id=calls.avl_id "
                    "WHERE calls.id=?", (cid,)).fetchone()
    if row:
        c.execute("DELETE FROM call_attendees WHERE call_id=?", (cid,))
        c.execute("DELETE FROM calls WHERE id=?", (cid,))
        c.commit()
        db.log(user["email"], "call:delete", f"{row['avl_name']} {row['call_date']}")
    c.close()
    return RedirectResponse(f"/calls?avl={row['avl_id'] if row else 0}", status_code=303)

# ---------------- executive dashboard ----------------
@app.get("/exec", response_class=HTMLResponse)
def exec_dash(request: Request, month: str = "", user=Depends(require_user)):
    c = db.conn()
    month = month or datetime.date.today().strftime("%Y-%m")
    changes = c.execute(
        "SELECT h.*, p.name AS product, a.name AS avl_name FROM status_history h "
        "JOIN products p ON p.id=h.product_id JOIN avls a ON a.id=h.avl_id "
        "WHERE substr(h.ts,1,7)=? ORDER BY h.ts", (month,)).fetchall()
    wins = [r for r in changes if r["new_status"] in db.LISTED_STATUSES]
    risks = [r for r in changes if r["new_status"] in ("No Interest", "N/A")]
    calls_month = c.execute(
        "SELECT a.name AS avl_name, COUNT(*) AS n FROM calls JOIN avls a ON a.id=calls.avl_id "
        "WHERE substr(call_date,1,7)=? GROUP BY a.name ORDER BY n DESC", (month,)).fetchall()
    n_calls = sum(r["n"] for r in calls_month)
    counts = c.execute(
        "SELECT l.status, COUNT(*) AS n FROM listings l "
        "JOIN products p ON p.id=l.product_id AND p.active=1 "
        "JOIN avls a ON a.id=l.avl_id AND a.active=1 GROUP BY l.status ORDER BY n DESC").fetchall()
    listed_by_avl = c.execute(
        "SELECT a.name, "
        "SUM(CASE WHEN l.status IN ('Listed','Listed, Conditional') THEN 1 ELSE 0 END) AS listed, "
        "SUM(CASE WHEN l.status='Listed, Conditional' THEN 1 ELSE 0 END) AS conditional, "
        "COUNT(*) AS total "
        "FROM listings l JOIN products p ON p.id=l.product_id AND p.active=1 "
        "JOIN avls a ON a.id=l.avl_id AND a.active=1 GROUP BY a.name ORDER BY listed DESC").fetchall()
    months = [r[0] for r in c.execute(
        "SELECT DISTINCT substr(ts,1,7) FROM status_history "
        "UNION SELECT DISTINCT substr(call_date,1,7) FROM calls ORDER BY 1 DESC")]
    c.close()
    return templates.TemplateResponse(request, "exec.html", {"user": user, "month": month,
        "months": months, "changes": changes, "wins": wins, "risks": risks,
        "calls_month": calls_month, "n_calls": n_calls, "counts": counts,
        "listed_by_avl": listed_by_avl})


@app.post("/products/{pid}/lifecycle")
def set_lifecycle(pid: int, request: Request, lifecycle: str = Form(...),
                  launch_status: str = Form(""), user=Depends(require_editor)):
    if lifecycle not in ("Roadmap", "Active", "EOL"):
        return RedirectResponse("/manage", status_code=303)
    c = db.conn()
    c.execute("UPDATE products SET lifecycle=?, launch_status=? WHERE id=?",
              (lifecycle, launch_status.strip(), pid))
    name = c.execute("SELECT name FROM products WHERE id=?", (pid,)).fetchone()["name"]
    c.commit(); c.close()
    db.log(user["email"], "product:lifecycle", f"{name} -> {lifecycle} ({launch_status})")
    return RedirectResponse("/manage", status_code=303)


# ---------------- team: people + trifecta assignments ----------------
@app.get("/team", response_class=HTMLResponse)
def team(request: Request, user=Depends(require_user)):
    c = db.conn()
    people = c.execute("SELECT * FROM people ORDER BY active DESC, name").fetchall()
    active_people = [p for p in people if p["active"]]
    load = {r["person_id"]: r["n"] for r in c.execute(
        "SELECT person_id, COUNT(*) AS n FROM assignments WHERE ended_at IS NULL GROUP BY person_id")}
    orgs = db.orgs_in_use(c)
    avls = c.execute("SELECT id, name FROM avls WHERE active=1 ORDER BY name").fetchall()
    products = c.execute("SELECT id, name FROM products WHERE active=1 ORDER BY category, name").fetchall()
    current = c.execute(
        "SELECT a.id, p.name AS person, a.role, av.name AS avl_name, pr.name AS product_name, a.started_at "
        "FROM assignments a JOIN people p ON p.id=a.person_id "
        "LEFT JOIN avls av ON av.id=a.avl_id LEFT JOIN products pr ON pr.id=a.product_id "
        "WHERE a.ended_at IS NULL ORDER BY av.name, pr.name, a.role").fetchall()
    past = c.execute(
        "SELECT p.name AS person, a.role, av.name AS avl_name, pr.name AS product_name, a.started_at, a.ended_at "
        "FROM assignments a JOIN people p ON p.id=a.person_id "
        "LEFT JOIN avls av ON av.id=a.avl_id LEFT JOIN products pr ON pr.id=a.product_id "
        "WHERE a.ended_at IS NOT NULL ORDER BY a.ended_at DESC LIMIT 50").fetchall()
    # coverage: per AVL, who fills each role
    cov = {}
    for r in current:
        if r["avl_name"]:
            cov.setdefault(r["avl_name"], {}).setdefault(r["role"], []).append(r["person"])
    c.close()
    return templates.TemplateResponse(request, "team.html", {"user": user, "people": people,
        "active_people": active_people, "load": load,
        "avls": avls, "products": products, "current": current, "past": past,
        "cov": cov, "roles": db.ROLES, "orgs": orgs})

def _org_value(org, org_other):
    """"Other" is a prompt to define the team, not a bucket to file people in."""
    o = (org or "").strip()
    return (org_other or "").strip() if o in ("", "Other") and org_other.strip() else o

@app.post("/team/person/add")
def add_person(request: Request, name: str = Form(...), email: str = Form(""),
               org: str = Form(""), org_other: str = Form(""), user=Depends(require_editor)):
    org = _org_value(org, org_other)
    nm = name.strip()
    if not nm:
        return RedirectResponse("/team", status_code=303)
    c = db.conn()
    row = c.execute("SELECT id, active FROM people WHERE lower(name)=lower(?)", (nm,)).fetchone()
    if row:
        # Re-adding someone who was retired brings them back rather than duplicating.
        c.execute("UPDATE people SET active=1, email=COALESCE(NULLIF(?,''), email), "
                  "org=COALESCE(NULLIF(?,''), org) WHERE id=?", (email.strip(), org.strip(), row["id"]))
        c.commit(); c.close()
        db.log(user["email"], "person:restore", nm)
        return RedirectResponse("/team", status_code=303)
    c.execute("INSERT INTO people(name, email, org) VALUES(?,?,?)",
              (nm, email.strip(), org.strip()))
    c.commit(); c.close()
    db.log(user["email"], "person:add", nm)
    return RedirectResponse("/team", status_code=303)

@app.post("/team/person/{person_id}/save")
def save_person(person_id: int, request: Request, name: str = Form(...), email: str = Form(""),
                org: str = Form(""), org_other: str = Form(""), user=Depends(require_editor)):
    org = _org_value(org, org_other)
    nm = name.strip()
    if not nm:
        return RedirectResponse("/team", status_code=303)
    c = db.conn()
    if c.execute("SELECT 1 FROM people WHERE lower(name)=lower(?) AND id<>?", (nm, person_id)).fetchone():
        c.close()
        return RedirectResponse("/team?err=dupname", status_code=303)
    c.execute("UPDATE people SET name=?, email=?, org=? WHERE id=?",
              (nm, email.strip(), org.strip(), person_id))
    db.refresh_all_rep_caches(c)
    c.commit(); c.close()
    db.log(user["email"], "person:save", nm)
    return RedirectResponse("/team", status_code=303)

@app.post("/team/person/{person_id}/toggle")
def toggle_person(person_id: int, request: Request, user=Depends(require_editor)):
    """Retiring someone ends their open assignments so they leave every picker."""
    c = db.conn()
    row = c.execute("SELECT name, active FROM people WHERE id=?", (person_id,)).fetchone()
    if not row:
        c.close()
        return RedirectResponse("/team", status_code=303)
    going_inactive = bool(row["active"])
    c.execute("UPDATE people SET active = 1 - active WHERE id=?", (person_id,))
    ended = 0
    if going_inactive:
        open_rows = c.execute("SELECT id, avl_id, product_id FROM assignments "
                              "WHERE person_id=? AND ended_at IS NULL", (person_id,)).fetchall()
        for a in open_rows:
            c.execute("UPDATE assignments SET ended_at=? WHERE id=?", (now(), a["id"]))
            db.refresh_rep_cache(c, avl_id=a["avl_id"], product_id=a["product_id"])
        ended = len(open_rows)
    c.commit(); c.close()
    db.log(user["email"], "person:toggle",
           f"{row['name']} -> {'retired' if going_inactive else 'active'}"
           + (f" ({ended} assignment(s) ended)" if ended else ""))
    return RedirectResponse("/team", status_code=303)

@app.post("/team/assign")
def assign(request: Request, person_id: int = Form(...), role: str = Form(...),
           avl_id: str = Form(""), product_id: str = Form(""), user=Depends(require_editor)):
    aid = int(avl_id) if avl_id else None
    pid = int(product_id) if product_id else None
    if role not in db.ROLES or (aid is None and pid is None):
        return RedirectResponse("/team", status_code=303)
    c = db.conn()
    if not c.execute("SELECT 1 FROM people WHERE id=? AND active=1", (person_id,)).fetchone():
        c.close()
        return RedirectResponse("/team?err=inactive", status_code=303)
    c.execute("INSERT INTO assignments(person_id, role, avl_id, product_id, started_at, added_by) "
              "VALUES(?,?,?,?,?,?)", (person_id, role, aid, pid, now(), user["email"]))
    pname = c.execute("SELECT name FROM people WHERE id=?", (person_id,)).fetchone()["name"]
    db.refresh_rep_cache(c, avl_id=aid, product_id=pid)
    c.commit(); c.close()
    db.log(user["email"], "assign:add", f"{pname} as {role} (avl={aid}, product={pid})")
    return RedirectResponse("/team", status_code=303)

@app.post("/team/assign/{assign_id}/end")
def end_assignment(assign_id: int, request: Request, user=Depends(require_editor)):
    c = db.conn()
    row = c.execute("SELECT avl_id, product_id FROM assignments WHERE id=?", (assign_id,)).fetchone()
    c.execute("UPDATE assignments SET ended_at=? WHERE id=? AND ended_at IS NULL", (now(), assign_id))
    if row:
        db.refresh_rep_cache(c, avl_id=row["avl_id"], product_id=row["product_id"])
    c.commit(); c.close()
    db.log(user["email"], "assign:end", f"assignment {assign_id}")
    return RedirectResponse("/team", status_code=303)
# ---- v6 features: appended to main.py ----
import io, csv, shutil, re as _re
from fastapi import UploadFile, File
from fastapi.responses import StreamingResponse, FileResponse

# Must be settable: on a hosted box the uploads belong on the same persistent
# volume as the database, not inside the checked-out code.
UPLOAD_DIR = os.environ.get("AVL_UPLOADS", os.path.join(os.path.dirname(BASE), "data_uploads"))
os.makedirs(UPLOAD_DIR, exist_ok=True)

def _csv_response(rows, header, fname):
    buf = io.StringIO()
    w = csv.writer(buf)
    w.writerow(header)
    for r in rows:
        w.writerow(r)
    buf.seek(0)
    return StreamingResponse(iter([buf.getvalue()]), media_type="text/csv",
                             headers={"Content-Disposition": f"attachment; filename={fname}"})

# ---------------- 1) per-cell notes ----------------
@app.post("/listing/note")
def listing_note(request: Request, product_id: int = Form(...), avl_id: int = Form(...),
                 note: str = Form(""), user=Depends(require_editor)):
    c = db.conn()
    c.execute("INSERT INTO listings(product_id, avl_id, status, note, updated_by, updated_at) "
              "VALUES(?,?, 'No Info', ?, ?, ?) "
              "ON CONFLICT(product_id, avl_id) DO UPDATE SET note=excluded.note, "
              "updated_by=excluded.updated_by, updated_at=excluded.updated_at",
              (product_id, avl_id, note.strip(), user["email"], now()))
    c.commit(); c.close()
    db.log(user["email"], "note", f"p{product_id}/a{avl_id}: {note[:80]}")
    return RedirectResponse("/", status_code=303)

# ---------------- 2) CSV exports ----------------
@app.get("/export/matrix.csv")
def export_matrix(request: Request, user=Depends(require_user)):
    c = db.conn()
    avls = c.execute("SELECT * FROM avls WHERE active=1 ORDER BY id").fetchall()
    products = c.execute("SELECT * FROM products WHERE active=1 ORDER BY category, id").fetchall()
    grid = {(r["product_id"], r["avl_id"]): r for r in c.execute("SELECT * FROM listings")}
    rows = []
    for p in products:
        row = [p["category"], p["name"], p["lifecycle"], p["launch_status"]]
        for a in avls:
            cell = grid.get((p["id"], a["id"]))
            row.append(cell["status"] if cell else "")
        rows.append(row)
    c.close()
    return _csv_response(rows, ["Category", "Product", "Lifecycle", "Launch"] + [a["name"] for a in avls],
                         "avl_matrix.csv")

@app.get("/export/history.csv")
def export_history(request: Request, user=Depends(require_user)):
    c = db.conn()
    rows = [(r["ts"], r["avl_name"], r["product"], r["old_status"] or "", r["new_status"], r["changed_by"])
            for r in c.execute("SELECT h.ts, h.old_status, h.new_status, h.changed_by, "
                               "p.name AS product, a.name AS avl_name FROM status_history h "
                               "JOIN products p ON p.id=h.product_id JOIN avls a ON a.id=h.avl_id "
                               "ORDER BY h.ts DESC")]
    c.close()
    return _csv_response(rows, ["When", "TPO AVL", "Product", "From", "To", "Changed by"], "avl_history.csv")

@app.get("/export/calls.csv")
def export_calls(request: Request, user=Depends(require_user)):
    c = db.conn()
    rows = [(r["call_date"], r["avl_name"], r["call_type"], r["qcells_attendees"], r["tpo_attendees"],
             r["topics"], r["outcomes"], r["owner_due"], r["created_by"])
            for r in c.execute("SELECT calls.*, a.name AS avl_name FROM calls "
                               "JOIN avls a ON a.id=calls.avl_id ORDER BY call_date DESC")]
    c.close()
    return _csv_response(rows, ["Date", "TPO", "Type", "Qcells attendees", "TPO attendees",
                                "Topics", "Outcomes", "Owner/Due", "Logged by"], "avl_calls.csv")

# ---------------- 3) action items ----------------
def _requirement_picker(c, avl_id=None):
    """avl -> product -> [requirement] for the checklist-link dropdowns."""
    q = ("SELECT ci.id, ci.workstream, ci.doc_category, ci.status, ci.product_id, ci.avl_id, "
         "p.name AS product FROM checklist_items ci JOIN products p ON p.id=ci.product_id ")
    args = []
    if avl_id:
        q += "WHERE ci.avl_id=? "
        args.append(avl_id)
    tree = {}
    for r in c.execute(q + "ORDER BY p.name, ci.sort_order, ci.id", args):
        tree.setdefault(str(r["avl_id"]), {}).setdefault(
            str(r["product_id"]), {"name": r["product"], "items": []})["items"].append(
            {"id": r["id"], "ws": r["workstream"], "cat": r["doc_category"], "status": r["status"]})
    return tree

@app.get("/actions/requirements")
def action_requirements(request: Request, avl_id: int = 0, product_id: int = 0,
                        user=Depends(require_user)):
    """Requirements selectable for one product x TPO: those already tracked, plus
    everything in the best-matching template that is not tracked yet."""
    if not avl_id or not product_id:
        return JSONResponse({"tracked": [], "template": [], "template_name": ""})
    c = db.conn()
    tracked = [{"value": f"ck:{r['id']}",
                "label": (f"{r['doc_category']}: " if r["doc_category"] else "") + r["workstream"],
                "status": r["status"]}
               for r in c.execute("SELECT * FROM checklist_items WHERE product_id=? AND avl_id=? "
                                  "ORDER BY sort_order, id", (product_id, avl_id))]
    have = {r["workstream"] for r in c.execute(
        "SELECT workstream FROM checklist_items WHERE product_id=? AND avl_id=?",
        (product_id, avl_id))}
    prow = c.execute("SELECT category FROM products WHERE id=?", (product_id,)).fetchone()
    tmpls = db.templates_for(c, prow["category"] if prow else "", avl_id)
    tmpl = tmpls[0] if tmpls else None
    template = []
    if tmpl:
        template = [{"value": f"tpl:{r['id']}",
                     "label": (f"{r['doc_category']}: " if r["doc_category"] else "") + r["workstream"],
                     "obligation": r["obligation"]}
                    for r in c.execute("SELECT * FROM workstream_template_items WHERE template_id=? "
                                       "ORDER BY sort_order, id", (tmpl["id"],))
                    if r["workstream"] not in have]
    c.close()
    return JSONResponse({"tracked": tracked, "template": template,
                         "template_name": tmpl["name"] if tmpl else ""})

@app.get("/actions", response_class=HTMLResponse)
def actions(request: Request, show: str = "open", avl: int = 0, owner: int = 0,
            listing: str = "", edit: int = 0, user=Depends(require_user)):
    c = db.conn()
    avls = c.execute("SELECT id, name FROM avls WHERE active=1 ORDER BY name").fetchall()
    people = c.execute("SELECT id, name, org FROM people WHERE active=1 ORDER BY name").fetchall()
    # Every active product is selectable, listed or not - an action is just as often
    # about getting a product onto an AVL as about a condition on one already there.
    products = c.execute("SELECT id, name, category FROM products WHERE active=1 "
                         "ORDER BY category, name").fetchall()
    q = ("SELECT actions.*, a.name AS avl_name, pe.name AS owner_name, "
         "pr.name AS product_name, ci.workstream, ci.doc_category, ci.status AS req_status, "
         "li.status AS listing_status FROM actions "
         "LEFT JOIN avls a ON a.id=actions.avl_id "
         "LEFT JOIN people pe ON pe.id=actions.owner_person_id "
         "LEFT JOIN checklist_items ci ON ci.id=actions.checklist_item_id "
         "LEFT JOIN products pr ON pr.id=actions.product_id "
         "LEFT JOIN listings li ON li.product_id=actions.product_id "
         "                     AND li.avl_id=actions.avl_id WHERE 1=1 ")
    params = []
    if show != "all":
        q += "AND actions.status='Open' "
    if avl:
        q += "AND actions.avl_id=? "
        params.append(avl)
    if owner:
        q += "AND actions.owner_person_id=? "
        params.append(owner)
    if listing == "listed":
        q += "AND li.status='Listed' "
    elif listing == "conditional":
        q += "AND li.status='Listed, Conditional' "
    elif listing == "unlisted":
        q += "AND actions.product_id IS NOT NULL AND COALESCE(li.status,'') NOT IN ('Listed','Listed, Conditional') "
    elif listing == "noproduct":
        q += "AND actions.product_id IS NULL "
    # actions and checklist_items both have due_date now, so qualify it.
    rows = c.execute(q + "ORDER BY CASE WHEN COALESCE(actions.due_date,'')='' THEN 1 ELSE 0 END, "
                     "actions.due_date, actions.id DESC", params).fetchall()
    today = datetime.date.today().isoformat()
    n_overdue = sum(1 for r in rows if r["status"] == "Open" and r["due_date"] and r["due_date"] < today)
    tree = _requirement_picker(c)
    listings = {f"{r['avl_id']}:{r['product_id']}": r["status"] for r in c.execute(
        "SELECT avl_id, product_id, status FROM listings")}
    c.close()
    return templates.TemplateResponse(request, "actions.html", {"user": user, "rows": rows,
        "avls": avls, "people": people, "show": show, "today": today, "sel_a": avl,
        "sel_o": owner, "edit": edit, "tree": tree, "priorities": db.ACTION_PRIORITIES,
        "n_overdue": n_overdue, "products": products, "listings": listings,
        "listed_statuses": db.LISTED_STATUSES, "sel_l": listing})

def _as_id(v):
    """0 and blank are the "none" sentinels the pickers use; never a real row."""
    return int(v) if str(v).strip().isdigit() and int(v) > 0 else None

def _resolve_requirement(c, requirement, avl_id, product_id, actor):
    """'ck:<id>' is an existing checklist row; 'tpl:<id>' starts tracking one.

    Attaching an untracked template requirement adds it to that product x TPO
    checklist, which is what "we owe them this" means in practice - and is the
    usual case when a listing is conditional on something nobody tracks yet.
    """
    kind, _, rid = (requirement or "").partition(":")
    if kind == "ck" and rid.isdigit():
        return int(rid)
    if kind != "tpl" or not rid.isdigit() or not (avl_id and product_id):
        return None
    t = c.execute("SELECT * FROM workstream_template_items WHERE id=?", (int(rid),)).fetchone()
    if not t:
        return None
    existing = c.execute("SELECT id FROM checklist_items WHERE product_id=? AND avl_id=? "
                         "AND workstream=?", (product_id, avl_id, t["workstream"])).fetchone()
    if existing:
        return existing["id"]
    nxt = c.execute("SELECT COALESCE(MAX(sort_order), 0) + 1 FROM checklist_items "
                    "WHERE product_id=? AND avl_id=?", (product_id, avl_id)).fetchone()[0]
    c.execute("INSERT INTO checklist_items(product_id, avl_id, doc_category, workstream, obligation, "
              "sort_order, template_id, updated_by, updated_at) VALUES(?,?,?,?,?,?,?,?,?)",
              (product_id, avl_id, t["doc_category"], t["workstream"], t["obligation"],
               max(nxt, t["sort_order"]), t["template_id"], actor, now()))
    return c.execute("SELECT last_insert_rowid() AS id").fetchone()["id"]

def _action_fields(c, avl_id, product_id, checklist_item_id, owner_person_id):
    """A linked requirement is authoritative for the AVL and product it belongs to."""
    aid, pid, cid = _as_id(avl_id), _as_id(product_id), _as_id(checklist_item_id)
    if cid:
        row = c.execute("SELECT avl_id, product_id FROM checklist_items WHERE id=?", (cid,)).fetchone()
        if row:
            aid, pid = row["avl_id"], row["product_id"]
        else:
            cid = None
    oid = _as_id(owner_person_id)
    oname = ""
    if oid:
        row = c.execute("SELECT name FROM people WHERE id=? AND active=1", (oid,)).fetchone()
        oname = row["name"] if row else ""
        oid = oid if row else None
    return aid, pid, cid, oid, oname

@app.post("/actions/add")
def add_action(request: Request, description: str = Form(...), avl_id: str = Form(""),
               product_id: str = Form(""), requirement: str = Form(""),
               owner_person_id: str = Form(""), owner_other: str = Form(""),
               due_date: str = Form(""), priority: str = Form("Normal"),
               call_id: str = Form(""), user=Depends(require_editor)):
    desc = description.strip()
    if not desc:
        return RedirectResponse("/actions", status_code=303)
    if priority not in db.ACTION_PRIORITIES:
        priority = "Normal"
    c = db.conn()
    cid0 = _resolve_requirement(c, requirement, _as_id(avl_id), _as_id(product_id), user["email"])
    aid, pid, cid, oid, oname = _action_fields(c, avl_id, product_id, cid0, owner_person_id)
    c.execute("INSERT INTO actions(avl_id, product_id, checklist_item_id, call_id, description, "
              "owner, owner_person_id, due_date, priority, created_by, created_at) "
              "VALUES(?,?,?,?,?,?,?,?,?,?,?)",
              (aid, pid, cid, int(call_id) if call_id.strip().isdigit() else None, desc,
               oname or owner_other.strip(), oid, due_date, priority, user["email"], now()))
    c.commit(); c.close()
    db.log(user["email"], "action:add", desc[:80], avl_id=aid, product_id=pid, entity="action")
    return RedirectResponse(f"/actions?avl={aid or 0}", status_code=303)

@app.post("/actions/{aid_}/save")
def save_action(aid_: int, request: Request, description: str = Form(...), avl_id: str = Form(""),
                product_id: str = Form(""), requirement: str = Form(""),
                owner_person_id: str = Form(""), owner_other: str = Form(""),
                due_date: str = Form(""), priority: str = Form("Normal"),
                user=Depends(require_editor)):
    desc = description.strip()
    if priority not in db.ACTION_PRIORITIES:
        priority = "Normal"
    c = db.conn()
    if not desc or not c.execute("SELECT 1 FROM actions WHERE id=?", (aid_,)).fetchone():
        c.close()
        return RedirectResponse("/actions", status_code=303)
    aid, pid, cid, oid, oname = _action_fields(c, avl_id, product_id, checklist_item_id,
                                               owner_person_id)
    prev = c.execute("SELECT owner FROM actions WHERE id=?", (aid_,)).fetchone()
    owner_txt = oname or owner_other.strip() or (prev["owner"] if prev and not oid else "")
    c.execute("UPDATE actions SET avl_id=?, product_id=?, checklist_item_id=?, description=?, "
              "owner=?, owner_person_id=?, due_date=?, priority=? WHERE id=?",
              (aid, pid, cid, desc, owner_txt, oid, due_date, priority, aid_))
    c.commit(); c.close()
    db.log(user["email"], "action:save", desc[:80], avl_id=aid, product_id=pid, entity="action")
    return RedirectResponse(f"/actions?avl={aid or 0}", status_code=303)

@app.post("/actions/{aid}/done")
def action_done(aid: int, request: Request, user=Depends(require_editor)):
    c = db.conn()
    c.execute("UPDATE actions SET status='Done', closed_at=? WHERE id=?", (now(), aid))
    row = c.execute("SELECT description FROM actions WHERE id=?", (aid,)).fetchone()
    c.commit(); c.close()
    db.log(user["email"], "action:done", (row["description"][:80] if row else str(aid)))
    return RedirectResponse("/actions", status_code=303)

@app.post("/actions/{aid}/reopen")
def action_reopen(aid: int, request: Request, user=Depends(require_editor)):
    c = db.conn()
    c.execute("UPDATE actions SET status='Open', closed_at=NULL WHERE id=?", (aid,))
    row = c.execute("SELECT description FROM actions WHERE id=?", (aid,)).fetchone()
    c.commit(); c.close()
    db.log(user["email"], "action:reopen", (row["description"][:80] if row else str(aid)))
    return RedirectResponse("/actions?show=all", status_code=303)

@app.post("/actions/{aid}/delete")
def action_delete(aid: int, request: Request, user=Depends(require_editor)):
    c = db.conn()
    row = c.execute("SELECT description FROM actions WHERE id=?", (aid,)).fetchone()
    c.execute("DELETE FROM actions WHERE id=?", (aid,))
    c.commit(); c.close()
    if row:
        db.log(user["email"], "action:delete", row["description"][:80])
    return RedirectResponse("/actions", status_code=303)

# ---------------- 4) PPTX exports ----------------
def _pptx_status_deck():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    FILLS = {"Listed": (0xC6,0xEF,0xCE), "In Review": (0xBD,0xD7,0xEE), "Execution": (0xA8,0xF0,0xEA),
             "Engagement": (0xFC,0xD9,0xC4), "Opportunity": (0xFC,0xE4,0xD6), "No Interest": (0xF8,0xCB,0xAD),
             "No Info": (0xFF,0xF2,0xCC), "N/A": (0xFF,0xFF,0xFF), "Pre-launch": (0xED,0xED,0xED)}
    NAVY = RGBColor(0x1F,0x38,0x64)
    pres = Presentation(); pres.slide_width = Inches(13.333); pres.slide_height = Inches(7.5)
    s = pres.slides.add_slide(pres.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.5))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = f"AVL Status - Current products ({datetime.date.today().strftime('%b %Y')})"
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = NAVY
    c = db.conn()
    avls = c.execute("SELECT * FROM avls WHERE active=1 ORDER BY id").fetchall()
    products = c.execute("SELECT * FROM products WHERE active=1 AND lifecycle != 'EOL' "
                         "ORDER BY category, id").fetchall()
    grid = {(x["product_id"], x["avl_id"]): x["status"] for x in c.execute("SELECT * FROM listings")}
    c.close()
    nrows, ncols = len(products) + 1, len(avls) + 1
    gf = s.shapes.add_table(nrows, ncols, Inches(0.4), Inches(0.75),
                            Inches(12.5), Inches(min(6.2, 0.34 * nrows)))
    tbl = gf.table
    tbl.columns[0].width = Inches(2.9)
    for i in range(1, ncols):
        tbl.columns[i].width = Inches((12.5 - 2.9) / len(avls))
    def cell(rr, cc, text, bold=False, fill=None, size=9):
        cl = tbl.cell(rr, cc); cl.text = text
        for para in cl.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(size); run.font.bold = bold
        if fill:
            cl.fill.solid(); cl.fill.fore_color.rgb = RGBColor(*fill)
    cell(0, 0, "Product", bold=True)
    for j, a in enumerate(avls, 1):
        cell(0, j, a["name"], bold=True, size=8)
    for i, p in enumerate(products, 1):
        label = p["name"] + (f" [{p['lifecycle']} {p['launch_status']}]" if p["lifecycle"] != "Active" else "")
        cell(i, 0, label, bold=True, size=8)
        for j, a in enumerate(avls, 1):
            st = grid.get((p["id"], a["id"]), "")
            cell(i, j, st, fill=FILLS.get(st), size=8)
    buf = io.BytesIO(); pres.save(buf); buf.seek(0)
    return buf

@app.get("/export/status.pptx")
def export_status_pptx(request: Request, user=Depends(require_user)):
    buf = _pptx_status_deck()
    return StreamingResponse(buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=AVL_Status_Snapshot.pptx"})

# ---------------- 5) dataroom checklists ----------------
def _checklist_atts(c, item_ids):
    """Attachments grouped by the checklist item they document."""
    out = {}
    if not item_ids:
        return out
    qs = ",".join("?" * len(item_ids))
    for r in c.execute(f"SELECT * FROM attachments WHERE kind='checklist' AND ref_id IN ({qs}) "
                       "ORDER BY id DESC", item_ids):
        out.setdefault(r["ref_id"], []).append(r)
    return out

@app.get("/dataroom", response_class=HTMLResponse)
def dataroom(request: Request, product: int = 0, avl: int = 0, only: str = "",
             user=Depends(require_user)):
    c = db.conn()
    products = c.execute("SELECT id, name, category FROM products WHERE active=1 "
                         "ORDER BY category, name").fetchall()
    avls = c.execute("SELECT id, name FROM avls WHERE active=1 ORDER BY name").fetchall()
    product = product or (products[0]["id"] if products else 0)
    avl = avl or (avls[0]["id"] if avls else 0)
    prow = c.execute("SELECT * FROM products WHERE id=?", (product,)).fetchone()
    category = prow["category"] if prow else ""
    q = "SELECT * FROM checklist_items WHERE product_id=? AND avl_id=? "
    params = [product, avl]
    if only in db.OBLIGATIONS:
        q += "AND obligation=? "
        params.append(only)
    elif only == "open":
        q += "AND status NOT IN ('Complete','TBD') "
    items = c.execute(q + "ORDER BY sort_order, id", params).fetchall()
    tmpls = db.templates_for(c, category, avl)
    atts = _checklist_atts(c, [i["id"] for i in items])
    today = datetime.date.today().isoformat()
    ie_reports = c.execute("SELECT id, name, reviewer, status FROM ie_reports "
                           "WHERE product_id=? AND active=1 ORDER BY id DESC",
                           (product,)).fetchall()
    ie_prog = db.ie_report_progress(c, [i["ie_report_id"] for i in items])

    # Group into the tracker's Document Categories, with a per-group rollup.
    # The filtered view must not distort the tracker, so the rollup is always
    # computed over the whole checklist.
    c_people = c.execute("SELECT id, name FROM people WHERE active=1 ORDER BY name").fetchall()
    all_items = items if not only else c.execute(
        "SELECT * FROM checklist_items WHERE product_id=? AND avl_id=? ORDER BY sort_order, id",
        (product, avl)).fetchall()
    all_atts = atts if not only else _checklist_atts(c, [i["id"] for i in all_items])

    groups, seen = [], {}
    for i in items:
        key = i["doc_category"] or "Other"
        if key not in seen:
            seen[key] = {"name": key, "items": []}
            groups.append(seen[key])
        seen[key]["items"].append(i)

    # Status tracker: one row per document category over the full checklist.
    track, tseen = [], {}
    for i in all_items:
        key = i["doc_category"] or "Other"
        if key not in tseen:
            tseen[key] = {"name": key, "n": 0, "req": 0, "done": 0, "wip": 0, "blocked": 0,
                          "open": 0, "files": 0, "nofile_req": 0}
            track.append(tseen[key])
        t = tseen[key]
        t["n"] += 1
        n_files = len(all_atts.get(i["id"], []))
        t["files"] += n_files
        required = i["obligation"] == "Required"
        if required:
            t["req"] += 1
        if i["status"] in db.CHECK_DONE:
            t["done"] += required
        elif i["status"] == "Blocked":
            t["blocked"] += 1
        elif i["status"] == "In Progress":
            t["wip"] += 1
        elif i["status"] == "Not Started":
            t["open"] += 1
        if required and not n_files:
            t["nofile_req"] += 1
    for t in track:
        t["pct"] = round(100 * t["done"] / t["req"]) if t["req"] else None
    totals = {k: sum(t[k] for t in track) for k in
              ("n", "req", "done", "wip", "blocked", "open", "files", "nofile_req")}
    totals["pct"] = round(100 * totals["done"] / totals["req"]) if totals["req"] else None
    total_req, total_done = totals["req"], totals["done"]
    c.close()
    return templates.TemplateResponse(request, "dataroom.html", {"user": user, "products": products,
        "avls": avls, "items": items, "groups": groups, "sel_p": product, "sel_a": avl,
        "check_statuses": db.CHECK_STATUSES, "obligations": db.OBLIGATIONS,
        "templates_": tmpls, "category": category, "atts": atts, "only": only,
        "total_req": total_req, "total_done": total_done, "track": track, "totals": totals,
        "today": today, "people": c_people, "ie_reports": ie_reports,
        "ie_prog": ie_prog, "looks_ie": db.looks_like_ie_requirement})

@app.post("/dataroom/seed")
def dataroom_seed(request: Request, product_id: int = Form(...), avl_id: int = Form(...),
                  template_id: int = Form(0), mode: str = Form("replace"),
                  user=Depends(require_editor)):
    """Apply a workstream template to one product x AVL checklist.

    mode=replace wipes rows that were never touched and re-seeds; mode=merge only
    adds workstreams that are missing, so in-flight status/notes survive.
    """
    c = db.conn()
    t = c.execute("SELECT * FROM workstream_templates WHERE id=?", (template_id,)).fetchone()
    if not t:
        c.close()
        return RedirectResponse(f"/dataroom?product={product_id}&avl={avl_id}", status_code=303)
    rows = c.execute("SELECT doc_category, workstream, obligation, sort_order "
                     "FROM workstream_template_items WHERE template_id=? "
                     "ORDER BY sort_order, id", (template_id,)).fetchall()
    if mode == "replace":
        # Keep anything a human has already worked on; only clear untouched rows.
        c.execute("DELETE FROM checklist_items WHERE product_id=? AND avl_id=? "
                  "AND status='Not Started' AND COALESCE(notes,'')='' AND COALESCE(pct,'')='' "
                  "AND COALESCE(eta,'')='' AND COALESCE(owner,'')='' "
                  "AND id NOT IN (SELECT ref_id FROM attachments WHERE kind='checklist')",
                  (product_id, avl_id))
    have = {r["workstream"] for r in c.execute(
        "SELECT workstream FROM checklist_items WHERE product_id=? AND avl_id=?",
        (product_id, avl_id))}
    added = 0
    for r in rows:
        if r["workstream"] in have:
            continue
        c.execute("INSERT INTO checklist_items(product_id, avl_id, doc_category, workstream, "
                  "obligation, sort_order, template_id, updated_by, updated_at) "
                  "VALUES(?,?,?,?,?,?,?,?,?)",
                  (product_id, avl_id, r["doc_category"], r["workstream"], r["obligation"],
                   r["sort_order"], template_id, user["email"], now()))
        added += 1
    c.commit(); c.close()
    db.log(user["email"], "dataroom:seed", f"applied '{t['name']}' ({mode}, +{added} requirements)",
           avl_id=avl_id, product_id=product_id, entity="checklist")
    return RedirectResponse(f"/dataroom?product={product_id}&avl={avl_id}", status_code=303)

@app.post("/dataroom/item/add")
def dataroom_item_add(request: Request, product_id: int = Form(...), avl_id: int = Form(...),
                      workstream: str = Form(...), doc_category: str = Form(""),
                      obligation: str = Form("Required"), user=Depends(require_editor)):
    """One-off requirement this TPO asked for that is not in any template."""
    ws = workstream.strip()
    if obligation not in db.OBLIGATIONS:
        obligation = "Required"
    if ws:
        c = db.conn()
        nxt = c.execute("SELECT COALESCE(MAX(sort_order), 0) + 1 FROM checklist_items "
                        "WHERE product_id=? AND avl_id=?", (product_id, avl_id)).fetchone()[0]
        c.execute("INSERT INTO checklist_items(product_id, avl_id, doc_category, workstream, "
                  "obligation, sort_order, updated_by, updated_at) VALUES(?,?,?,?,?,?,?,?)",
                  (product_id, avl_id, doc_category.strip(), ws, obligation, nxt,
                   user["email"], now()))
        c.commit(); c.close()
        db.log(user["email"], "dataroom:item:add", ws,
               avl_id=avl_id, product_id=product_id, entity="checklist")
    return RedirectResponse(f"/dataroom?product={product_id}&avl={avl_id}", status_code=303)

@app.post("/dataroom/item/{iid}/delete")
def dataroom_item_delete(iid: int, request: Request, user=Depends(require_editor)):
    c = db.conn()
    row = c.execute("SELECT * FROM checklist_items WHERE id=?", (iid,)).fetchone()
    if not row:
        c.close()
        return RedirectResponse("/dataroom", status_code=303)
    n_files = c.execute("SELECT COUNT(*) FROM attachments WHERE kind='checklist' AND ref_id=?",
                        (iid,)).fetchone()[0]
    if n_files:
        # Files would be orphaned; make the user move them first.
        c.close()
        return RedirectResponse(f"/dataroom?product={row['product_id']}&avl={row['avl_id']}"
                                f"&err=files", status_code=303)
    c.execute("DELETE FROM checklist_items WHERE id=?", (iid,))
    c.commit(); c.close()
    db.log(user["email"], "dataroom:item:delete", row["workstream"],
           avl_id=row["avl_id"], product_id=row["product_id"], entity="checklist")
    return RedirectResponse(f"/dataroom?product={row['product_id']}&avl={row['avl_id']}",
                            status_code=303)

@app.post("/dataroom/item/{iid}/update")
def dataroom_update(iid: int, request: Request, status: str = Form(...), pct: str = Form(""),
                    eta: str = Form(""), owner: str = Form(""), notes: str = Form(""),
                    due_date: str = Form(""), user=Depends(require_editor)):
    c = db.conn()
    c.execute("UPDATE checklist_items SET status=?, pct=?, eta=?, owner=?, notes=?, due_date=?, "
              "updated_by=?, updated_at=? WHERE id=?",
              (status, pct.strip(), eta.strip(), owner.strip(), notes.strip(), due_date,
               user["email"], now(), iid))
    row = c.execute("SELECT product_id, avl_id, workstream FROM checklist_items WHERE id=?", (iid,)).fetchone()
    c.commit(); c.close()
    db.log(user["email"], "dataroom:update", f"{row['workstream']} -> {status}",
           avl_id=row["avl_id"], product_id=row["product_id"],
           entity="checklist", entity_id=iid)
    return RedirectResponse(f"/dataroom?product={row['product_id']}&avl={row['avl_id']}", status_code=303)

@app.post("/dataroom/item/{iid}/link")
def dataroom_link_ie(iid: int, request: Request, ie_report_id: str = Form(""),
                     user=Depends(require_editor)):
    """Point a dataroom requirement at the IE review that answers it."""
    c = db.conn()
    row = c.execute("SELECT * FROM checklist_items WHERE id=?", (iid,)).fetchone()
    if not row:
        c.close()
        return RedirectResponse("/dataroom", status_code=303)
    rid = _as_id(ie_report_id)
    if rid:
        # Only a review of this product can answer this product's requirement.
        ok = c.execute("SELECT name FROM ie_reports WHERE id=? AND product_id=? AND active=1",
                       (rid, row["product_id"])).fetchone()
        if not ok:
            c.close()
            return RedirectResponse(f"/dataroom?product={row['product_id']}&avl={row['avl_id']}"
                                    "&err=iereport", status_code=303)
    c.execute("UPDATE checklist_items SET ie_report_id=?, updated_by=?, updated_at=? WHERE id=?",
              (rid, user["email"], now(), iid))
    c.commit(); c.close()
    db.log(user["email"], "dataroom:ie-link",
           f"{row['workstream']} {'-> IE report #' + str(rid) if rid else 'unlinked'}",
           avl_id=row["avl_id"], product_id=row["product_id"], entity="checklist", entity_id=iid)
    return RedirectResponse(f"/dataroom?product={row['product_id']}&avl={row['avl_id']}",
                            status_code=303)

@app.post("/dataroom/bulk")
def dataroom_bulk(request: Request, product_id: int = Form(...), avl_id: int = Form(...),
                  doc_category: str = Form(""), owner: str = Form(""),
                  due_date: str = Form(""), status: str = Form(""),
                  user=Depends(require_editor)):
    """Set owner / due / status across one document category at once."""
    sets, args = [], []
    if owner.strip():
        sets.append("owner=?"); args.append(owner.strip())
    if due_date:
        sets.append("due_date=?"); args.append(due_date)
    if status in db.CHECK_STATUSES:
        sets.append("status=?"); args.append(status)
    if sets:
        c = db.conn()
        c.execute(f"UPDATE checklist_items SET {', '.join(sets)}, updated_by=?, updated_at=? "
                  "WHERE product_id=? AND avl_id=? AND COALESCE(doc_category,'')=?",
                  args + [user["email"], now(), product_id, avl_id, doc_category])
        n = c.total_changes
        c.commit(); c.close()
        db.log(user["email"], "dataroom:bulk", f"{doc_category or 'Other'}: {n} requirement(s)",
               avl_id=avl_id, product_id=product_id, entity="checklist")
    return RedirectResponse(f"/dataroom?product={product_id}&avl={avl_id}", status_code=303)

@app.get("/dataroom/queue", response_class=HTMLResponse)
def dataroom_queue(request: Request, owner: str = "", status: str = "", avl: int = 0,
                   overdue: int = 0, user=Depends(require_user)):
    """Requirements across every product x TPO - the view a person works from."""
    c = db.conn()
    today = datetime.date.today().isoformat()
    q = ("SELECT ci.*, p.name AS product, a.name AS avl_name FROM checklist_items ci "
         "JOIN products p ON p.id=ci.product_id AND p.active=1 "
         "JOIN avls a ON a.id=ci.avl_id AND a.active=1 WHERE 1=1 ")
    args = []
    if owner:
        q += "AND ci.owner=? "; args.append(owner)
    if status in db.CHECK_STATUSES:
        q += "AND ci.status=? "; args.append(status)
    elif status == "open":
        q += "AND ci.status NOT IN ('Complete','Submitted','Accepted') "
    if avl:
        q += "AND ci.avl_id=? "; args.append(avl)
    if overdue:
        q += "AND COALESCE(ci.due_date,'')<>'' AND ci.due_date<? " \
             "AND ci.status NOT IN ('Complete','Submitted','Accepted') "
        args.append(today)
    rows = c.execute(q + "ORDER BY CASE WHEN COALESCE(ci.due_date,'')='' THEN 1 ELSE 0 END, "
                     "ci.due_date, a.name, p.name, ci.sort_order LIMIT 400", args).fetchall()
    owners = [r[0] for r in c.execute("SELECT DISTINCT owner FROM checklist_items "
                                      "WHERE COALESCE(owner,'')<>'' ORDER BY owner")]
    avls = c.execute("SELECT id, name FROM avls WHERE active=1 ORDER BY name").fetchall()
    n_overdue = c.execute("SELECT COUNT(*) FROM checklist_items WHERE COALESCE(due_date,'')<>'' "
                          "AND due_date<? AND status NOT IN ('Complete','Submitted','Accepted')",
                          (today,)).fetchone()[0]
    c.close()
    return templates.TemplateResponse(request, "dataroom_queue.html", {"user": user, "rows": rows,
        "owners": owners, "avls": avls, "statuses": db.CHECK_STATUSES, "today": today,
        "sel": {"owner": owner, "status": status, "avl": avl, "overdue": overdue},
        "n_overdue": n_overdue})

# ---------------- 5b) workstream template management ----------------
@app.get("/workstreams", response_class=HTMLResponse)
def workstreams(request: Request, t: int = 0, user=Depends(require_user)):
    c = db.conn()
    avls = c.execute("SELECT id, name FROM avls WHERE active=1 ORDER BY name").fetchall()
    tmpls = c.execute("SELECT t.*, a.name AS avl_name, "
                      "(SELECT COUNT(*) FROM workstream_template_items i WHERE i.template_id=t.id) AS n_items "
                      "FROM workstream_templates t LEFT JOIN avls a ON a.id=t.avl_id "
                      "ORDER BY t.active DESC, t.category, t.name").fetchall()
    sel = t or (tmpls[0]["id"] if tmpls else 0)
    cur = c.execute("SELECT t.*, a.name AS avl_name FROM workstream_templates t "
                    "LEFT JOIN avls a ON a.id=t.avl_id WHERE t.id=?", (sel,)).fetchone()
    items = c.execute("SELECT * FROM workstream_template_items WHERE template_id=? "
                      "ORDER BY sort_order, id", (sel,)).fetchall() if cur else []
    doc_cats = [r[0] for r in c.execute(
        "SELECT DISTINCT doc_category FROM workstream_template_items "
        "WHERE doc_category<>'' ORDER BY doc_category")]
    # Where is this template already in play?
    history = c.execute("SELECT * FROM template_revisions WHERE template_id=? "
                        "ORDER BY id DESC LIMIT 15", (sel,)).fetchall() if cur else []
    usage = c.execute("SELECT p.name AS product, a.name AS avl_name, COUNT(*) AS n "
                      "FROM checklist_items ci JOIN products p ON p.id=ci.product_id "
                      "JOIN avls a ON a.id=ci.avl_id WHERE ci.template_id=? "
                      "GROUP BY ci.product_id, ci.avl_id ORDER BY p.name", (sel,)).fetchall() if cur else []
    c.close()
    return templates.TemplateResponse(request, "workstreams.html", {"user": user, "tmpls": tmpls,
        "cur": cur, "items": items, "avls": avls, "categories": db.CATEGORIES, "usage": usage,
        "obligations": db.OBLIGATIONS, "doc_cats": doc_cats, "history": history})

@app.post("/workstreams/add")
def ws_template_add(request: Request, name: str = Form(...), category: str = Form(""),
                    avl_id: str = Form(""), notes: str = Form(""), source_url: str = Form(""),
                    copy_from: str = Form(""), user=Depends(require_editor)):
    nm = name.strip()
    if not nm:
        return RedirectResponse("/workstreams", status_code=303)
    if category and category not in db.CATEGORIES:
        category = ""
    aid = int(avl_id) if avl_id.strip().isdigit() else None
    c = db.conn()
    if c.execute("SELECT 1 FROM workstream_templates WHERE name=?", (nm,)).fetchone():
        c.close()
        return RedirectResponse("/workstreams?err=dup", status_code=303)
    c.execute("INSERT INTO workstream_templates(name, category, avl_id, notes, source_url, "
              "created_by, created_at) VALUES(?,?,?,?,?,?,?)",
              (nm, category, aid, notes.strip(), source_url.strip(), user["email"], now()))
    tid = c.execute("SELECT id FROM workstream_templates WHERE name=?", (nm,)).fetchone()["id"]
    if copy_from.strip().isdigit():
        c.execute("INSERT INTO workstream_template_items(template_id, doc_category, workstream, "
                  "obligation, description, sort_order) "
                  "SELECT ?, doc_category, workstream, obligation, description, sort_order "
                  "FROM workstream_template_items WHERE template_id=?", (tid, int(copy_from)))
    c.commit(); c.close()
    db.log(user["email"], "workstream:template:add", f"{nm} [{category or 'any'}/{aid or 'any AVL'}]")
    return RedirectResponse(f"/workstreams?t={tid}", status_code=303)

@app.post("/workstreams/{tid}/edit")
def ws_template_edit(tid: int, request: Request, name: str = Form(...), category: str = Form(""),
                     avl_id: str = Form(""), notes: str = Form(""), source_url: str = Form(""),
                     user=Depends(require_editor)):
    nm = name.strip()
    if not nm:
        return RedirectResponse(f"/workstreams?t={tid}", status_code=303)
    if category and category not in db.CATEGORIES:
        category = ""
    aid = int(avl_id) if avl_id.strip().isdigit() else None
    c = db.conn()
    if c.execute("SELECT 1 FROM workstream_templates WHERE name=? AND id<>?", (nm, tid)).fetchone():
        c.close()
        return RedirectResponse(f"/workstreams?t={tid}&err=dup", status_code=303)
    db.record_revision(c, tid, "scope", f"renamed/rescoped to '{nm}'", user["email"])
    c.execute("UPDATE workstream_templates SET name=?, category=?, avl_id=?, notes=?, source_url=? "
              "WHERE id=?", (nm, category, aid, notes.strip(), source_url.strip(), tid))
    c.commit(); c.close()
    db.log(user["email"], "workstream:template:edit", f"{nm} [{category or 'any'}/{aid or 'any AVL'}]")
    return RedirectResponse(f"/workstreams?t={tid}", status_code=303)

@app.post("/workstreams/{tid}/toggle")
def ws_template_toggle(tid: int, request: Request, user=Depends(require_editor)):
    c = db.conn()
    db.record_revision(c, tid, "retire/restore", "", user["email"])
    c.execute("UPDATE workstream_templates SET active = 1 - active WHERE id=?", (tid,))
    row = c.execute("SELECT name, active FROM workstream_templates WHERE id=?", (tid,)).fetchone()
    c.commit(); c.close()
    db.log(user["email"], "workstream:template:toggle",
           f"{row['name']} -> {'active' if row['active'] else 'retired'}")
    return RedirectResponse(f"/workstreams?t={tid}", status_code=303)

@app.post("/workstreams/{tid}/delete")
def ws_template_delete(tid: int, request: Request, user=Depends(require_admin)):
    c = db.conn()
    row = c.execute("SELECT name FROM workstream_templates WHERE id=?", (tid,)).fetchone()
    c.execute("DELETE FROM workstream_templates WHERE id=?", (tid,))
    c.commit(); c.close()
    if row:
        db.log(user["email"], "workstream:template:delete", row["name"])
    return RedirectResponse("/workstreams", status_code=303)

@app.post("/workstreams/undo/{rev_id}")
def ws_undo(rev_id: int, request: Request, user=Depends(require_editor)):
    """Roll a template back to how it was before one recorded change."""
    c = db.conn()
    res = db.restore_revision(c, rev_id, user["email"])
    if not res:
        c.close()
        return RedirectResponse("/workstreams?err=gone", status_code=303)
    tid, action = res
    c.commit(); c.close()
    db.log(user["email"], "workstream:undo", f"t{tid}: reverted '{action}'")
    return RedirectResponse(f"/workstreams?t={tid}&undone={action}", status_code=303)

@app.post("/workstreams/{tid}/items/add")
def ws_item_add(tid: int, request: Request, workstream: str = Form(...),
                doc_category: str = Form(""), obligation: str = Form("Required"),
                description: str = Form(""), user=Depends(require_editor)):
    ws = workstream.strip()
    if obligation not in db.OBLIGATIONS:
        obligation = "Required"
    if ws:
        c = db.conn()
        nxt = c.execute("SELECT COALESCE(MAX(sort_order), 0) + 1 FROM workstream_template_items "
                        "WHERE template_id=?", (tid,)).fetchone()[0]
        db.record_revision(c, tid, "add requirement", ws, user["email"])
        c.execute("INSERT INTO workstream_template_items(template_id, doc_category, workstream, "
                  "obligation, description, sort_order) VALUES(?,?,?,?,?,?)",
                  (tid, doc_category.strip(), ws, obligation, description.strip(), nxt))
        c.commit(); c.close()
        db.log(user["email"], "workstream:item:add", f"t{tid} {ws}")
    return RedirectResponse(f"/workstreams?t={tid}", status_code=303)

@app.post("/workstreams/items/{iid}/save")
def ws_item_save(iid: int, request: Request, workstream: str = Form(...),
                 doc_category: str = Form(""), obligation: str = Form("Required"),
                 description: str = Form(""), sort_order: int = Form(0),
                 user=Depends(require_editor)):
    if obligation not in db.OBLIGATIONS:
        obligation = "Required"
    c = db.conn()
    row0 = c.execute("SELECT template_id, workstream FROM workstream_template_items WHERE id=?",
                     (iid,)).fetchone()
    if not row0:
        c.close()
        return RedirectResponse("/workstreams", status_code=303)
    tid = row0["template_id"]
    db.record_revision(c, tid, "edit requirement", row0["workstream"], user["email"])
    c.execute("UPDATE workstream_template_items SET doc_category=?, workstream=?, obligation=?, "
              "description=?, sort_order=? WHERE id=?",
              (doc_category.strip(), workstream.strip(), obligation, description.strip(),
               sort_order, iid))
    c.commit(); c.close()
    db.log(user["email"], "workstream:item:save", workstream.strip())
    return RedirectResponse(f"/workstreams?t={tid}", status_code=303)

@app.post("/workstreams/items/{iid}/delete")
def ws_item_delete(iid: int, request: Request, user=Depends(require_editor)):
    c = db.conn()
    row = c.execute("SELECT * FROM workstream_template_items WHERE id=?", (iid,)).fetchone()
    if not row:
        c.close()
        return RedirectResponse("/workstreams", status_code=303)
    db.record_revision(c, row["template_id"], "delete requirement", row["workstream"], user["email"])
    c.execute("DELETE FROM workstream_template_items WHERE id=?", (iid,))
    c.commit(); c.close()
    db.log(user["email"], "workstream:item:delete", row["workstream"])
    return RedirectResponse(f"/workstreams?t={row['template_id']}", status_code=303)

# ---------------- 6) weekly digest ----------------
def _digest_data(days=7):
    since = (datetime.date.today() - datetime.timedelta(days=days)).isoformat()
    c = db.conn()
    changes = c.execute("SELECT h.ts, h.old_status, h.new_status, p.name AS product, a.name AS avl_name "
                        "FROM status_history h JOIN products p ON p.id=h.product_id "
                        "JOIN avls a ON a.id=h.avl_id WHERE h.ts>=? ORDER BY h.ts DESC", (since,)).fetchall()
    calls = c.execute("SELECT calls.*, a.name AS avl_name FROM calls JOIN avls a ON a.id=calls.avl_id "
                      "WHERE call_date>=? ORDER BY call_date DESC", (since,)).fetchall()
    today = datetime.date.today().isoformat()
    overdue = c.execute("SELECT actions.*, a.name AS avl_name FROM actions "
                        "LEFT JOIN avls a ON a.id=actions.avl_id "
                        "WHERE status='Open' AND due_date != '' AND due_date < ? ORDER BY due_date",
                        (today,)).fetchall()
    c.close()
    return {"since": since, "changes": changes, "calls": calls, "overdue": overdue}

@app.get("/digest", response_class=HTMLResponse)
def digest_preview(request: Request, user=Depends(require_user)):
    d = _digest_data()
    return templates.TemplateResponse(request, "digest.html", {"user": user, **d})

# ---------------- 7-8) admin: roles + backup ----------------
@app.get("/admin", response_class=HTMLResponse)
def admin(request: Request, user=Depends(require_admin)):
    c = db.conn()
    users = c.execute("SELECT * FROM users ORDER BY role, email").fetchall()
    c.close()
    snaps = db.list_snapshots()[:10]
    return templates.TemplateResponse(request, "admin.html", {"user": user, "users": users,
        "snapshots": snaps, "backup_hours": db.AUTO_BACKUP_HOURS,
        "backup_keep": db.AUTO_BACKUP_KEEP, "backup_dir": db.BACKUP_DIR,
        "backup_token_set": bool(_auth.BACKUP_TOKEN), "n_uploads": _upload_count(),
        "auth_mode": AUTH_MODE, "allowlist": _auth.REQUIRE_KNOWN_USER,
        "admin_emails": sorted(_auth.ADMIN_EMAILS),
        "n_admins": sum(1 for u in users if u["role"] == "admin"),
        "warnings": _auth.startup_warnings(), "domain": ALLOWED_DOMAIN})

@app.post("/admin/user/add")
def admin_add_user(request: Request, email: str = Form(...), name: str = Form(""),
                   role: str = Form("editor"), user=Depends(require_admin)):
    """Pre-create someone so their role is right before they ever sign in."""
    em = email.strip().lower()
    if role not in ("viewer", "editor", "admin"):
        role = "editor"
    if not domain_ok(em):
        return RedirectResponse("/admin?err=domain", status_code=303)
    c = db.conn()
    if c.execute("SELECT 1 FROM users WHERE lower(email)=?", (em,)).fetchone():
        c.close()
        return RedirectResponse("/admin?err=exists", status_code=303)
    c.execute("INSERT INTO users(email, name, role) VALUES(?,?,?)",
              (em, name.strip() or em.split("@")[0], role))
    c.commit(); c.close()
    db.log(user["email"], "user:add", f"{em} as {role}")
    return RedirectResponse("/admin?ok=added", status_code=303)

@app.post("/admin/user/remove")
def admin_remove_user(request: Request, email: str = Form(...), user=Depends(require_admin)):
    em = email.strip().lower()
    if em == user["email"].lower():
        return RedirectResponse("/admin?err=self", status_code=303)
    c = db.conn()
    row = c.execute("SELECT role FROM users WHERE lower(email)=?", (em,)).fetchone()
    if not row:
        c.close()
        return RedirectResponse("/admin", status_code=303)
    if row["role"] == "admin" and \
       c.execute("SELECT COUNT(*) FROM users WHERE role='admin'").fetchone()[0] <= 1:
        c.close()
        return RedirectResponse("/admin?err=lastadmin", status_code=303)
    c.execute("DELETE FROM users WHERE lower(email)=?", (em,))
    c.commit(); c.close()
    db.log(user["email"], "user:remove", em)
    return RedirectResponse("/admin?ok=removed", status_code=303)

@app.post("/admin/role")
def set_role(request: Request, email: str = Form(...), role: str = Form(...),
             user=Depends(require_admin)):
    if role not in ("viewer", "editor", "admin"):
        return RedirectResponse("/admin", status_code=303)
    c = db.conn()
    cur = c.execute("SELECT role FROM users WHERE email=?", (email,)).fetchone()
    if cur and cur["role"] == "admin" and role != "admin" and \
       c.execute("SELECT COUNT(*) FROM users WHERE role='admin'").fetchone()[0] <= 1:
        c.close()
        return RedirectResponse("/admin?err=lastadmin", status_code=303)
    c.execute("UPDATE users SET role=? WHERE email=?", (role, email))
    c.commit(); c.close()
    db.log(user["email"], "role", f"{email} -> {role}")
    return RedirectResponse("/admin", status_code=303)

def _upload_count():
    try:
        return sum(1 for n in os.listdir(UPLOAD_DIR)
                   if os.path.isfile(os.path.join(UPLOAD_DIR, n)))
    except OSError:
        return 0

@app.get("/admin/backup/{name}")
def admin_backup_download(name: str, request: Request, user=Depends(require_admin)):
    """Download one of the automatic snapshots."""
    safe = os.path.basename(name)
    path = os.path.join(db.BACKUP_DIR, safe)
    if not (safe.startswith("avl_") and safe.endswith(".db") and os.path.exists(path)):
        return RedirectResponse("/admin", status_code=303)
    return FileResponse(path, filename=safe)

@app.get("/backup/{token}")
def backup_pull(token: str, request: Request, full: int = 0):
    """Unattended backup for a scheduler outside this box.

    Guarded by a long shared token rather than a session, because the point is to
    be fetchable by cron from somewhere else. Off by default: with no BACKUP_TOKEN
    set this route does not exist as far as a caller can tell.

    full=1 returns database *and* documents. The database alone is not a complete
    backup - it stores the path of each attachment, not its contents.
    """
    if not _auth.backup_token_ok(token):
        raise HTTPException(status_code=404)
    stamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    from starlette.background import BackgroundTask
    import tempfile
    if not full:
        fd, tmp = tempfile.mkstemp(prefix="pull_", suffix=".db"); os.close(fd)
        db.snapshot(tmp)
        db.log("backup-token", "backup:pull", f"database only ({os.path.getsize(tmp)} bytes)")
        return FileResponse(tmp, filename=f"avl_{stamp}.db", media_type="application/octet-stream",
                            background=BackgroundTask(lambda: os.path.exists(tmp) and os.remove(tmp)))
    import tarfile
    fd, tmpdb = tempfile.mkstemp(prefix="pull_", suffix=".db"); os.close(fd)
    db.snapshot(tmpdb)
    fd, tar = tempfile.mkstemp(prefix="pull_", suffix=".tar.gz"); os.close(fd)
    with tarfile.open(tar, "w:gz") as t:
        t.add(tmpdb, arcname=f"avl_{stamp}/avl.db")
        if os.path.isdir(UPLOAD_DIR):
            t.add(UPLOAD_DIR, arcname=f"avl_{stamp}/data_uploads")
    os.remove(tmpdb)
    db.log("backup-token", "backup:pull", f"full archive ({os.path.getsize(tar)} bytes)")
    return FileResponse(tar, filename=f"avl_full_{stamp}.tar.gz", media_type="application/gzip",
                        background=BackgroundTask(lambda: os.path.exists(tar) and os.remove(tar)))

@app.get("/admin/backup")
def backup(request: Request, user=Depends(require_admin)):
    """Download a consistent snapshot, leaving nothing behind.

    A plain file copy is not safe in WAL mode - recent commits live in the -wal
    file and would be missing - so this uses SQLite's backup API. The temp file
    is deleted once the response has been sent, rather than accumulating in the
    uploads directory on a 1 GB disk.
    """
    import sqlite3, tempfile
    from starlette.background import BackgroundTask
    stamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    fd, tmp = tempfile.mkstemp(prefix=f"avl_backup_{stamp}_", suffix=".db")
    os.close(fd)
    src = sqlite3.connect(db.DB_PATH)
    dst = sqlite3.connect(tmp)
    with dst:
        src.backup(dst)
    src.close(); dst.close()
    db.log(user["email"], "backup", f"downloaded ({os.path.getsize(tmp)} bytes)")
    return FileResponse(tmp, filename=f"avl_backup_{stamp}.db",
                        background=BackgroundTask(lambda: os.path.exists(tmp) and os.remove(tmp)))

# ---------------- 9) attachments ----------------
ATT_KINDS = ("product", "avl", "call", "checklist", "ie_item", "ie_section")

def _checklist_targets(c):
    """Every product x AVL x workstream row a file can be filed against."""
    return c.execute(
        "SELECT ci.id, ci.workstream, ci.status, ci.doc_category, ci.obligation, "
        "ci.product_id, ci.avl_id, p.name AS product, a.name AS avl_name "
        "FROM checklist_items ci JOIN products p ON p.id=ci.product_id "
        "JOIN avls a ON a.id=ci.avl_id "
        "ORDER BY a.name, p.name, ci.sort_order, ci.id").fetchall()

def _ie_targets(c):
    """Every IE report item a file can be filed against, plus a picker tree."""
    rows = c.execute(
        "SELECT i.id, i.item_id, i.sub_section, i.review_item, i.status, i.section_id, "
        "s.title AS section, s.sort_order AS s_ord, r.id AS report_id, r.name AS report, "
        "r.reviewer, p.name AS product "
        "FROM ie_report_items i JOIN ie_report_sections s ON s.id=i.section_id "
        "JOIN ie_reports r ON r.id=i.report_id JOIN products p ON p.id=r.product_id "
        "WHERE r.active=1 ORDER BY r.id, s.sort_order, i.sort_order, i.id").fetchall()
    tree, reports, sections = {}, [], {}
    for r in rows:
        rk, sk = str(r["report_id"]), str(r["section_id"])
        if rk not in tree:
            tree[rk] = {}
            reports.append([r["report_id"], f"{r['product']} - {r['reviewer']}"])
        if sk not in tree[rk]:
            tree[rk][sk] = {"name": r["section"], "items": []}
            sections[sk] = r["section"]
        tree[rk][sk]["items"].append(
            {"id": r["id"], "label": f"{r['item_id']} {r['sub_section']}".strip(),
             "status": r["status"]})
    return rows, {"tree": tree, "reports": reports}

def _files_context(c):
    products = c.execute("SELECT id, name FROM products ORDER BY name").fetchall()
    avls = c.execute("SELECT id, name FROM avls ORDER BY name").fetchall()
    calls_ = c.execute("SELECT calls.id, call_date, a.name AS avl_name FROM calls "
                       "JOIN avls a ON a.id=calls.avl_id ORDER BY call_date DESC LIMIT 100").fetchall()
    checks = _checklist_targets(c)
    ie_rows, ie_picker = _ie_targets(c)
    names = {"ie_item": {r["id"]: f"{r['product']} / {r['reviewer']} - {r['section']} / "
                                  f"{r['item_id']} {r['sub_section']}".strip() for r in ie_rows},
             "product": {p["id"]: p["name"] for p in products},
             "avl": {a["id"]: a["name"] for a in avls},
             "call": {r["id"]: f"{r['avl_name']} {r['call_date']}" for r in calls_},
             "checklist": {r["id"]: f"{r['avl_name']} / {r['product']} - "
                                    f"{r['doc_category'] + ': ' if r['doc_category'] else ''}"
                                    f"{r['workstream']}" for r in checks}}
    # Cascading picker data for the upload form: avl -> product -> [checklist rows]
    tree = {}
    for r in checks:
        tree.setdefault(str(r["avl_id"]), {}).setdefault(str(r["product_id"]), []).append(
            {"id": r["id"], "ws": r["workstream"], "status": r["status"],
             "cat": r["doc_category"]})
    picker = {"tree": tree,
              "avls": [[a["id"], a["name"]] for a in avls if str(a["id"]) in tree],
              "prods": {str(p["id"]): p["name"] for p in products}}
    return products, avls, calls_, checks, names, picker, ie_rows, ie_picker

@app.get("/files", response_class=HTMLResponse)
def files(request: Request, kind: str = "", ref: int = 0, user=Depends(require_user)):
    c = db.conn()
    products, avls, calls_, checks, names, picker, ie_rows, ie_picker = _files_context(c)
    q = "SELECT * FROM attachments"
    params = []
    if kind in ATT_KINDS:
        q += " WHERE kind=?"
        params.append(kind)
        if ref:
            q += " AND ref_id=?"
            params.append(ref)
    atts = c.execute(q + " ORDER BY id DESC", params).fetchall()
    # How many files sit under each bundle target, so the download side can say
    # up front whether there is anything to fetch.
    counts = {"product": {}, "avl": {}, "call": {}, "checklist": {}}
    for r in c.execute("SELECT kind, ref_id, COUNT(*) n FROM attachments GROUP BY kind, ref_id"):
        counts.setdefault(r["kind"], {})[r["ref_id"]] = r["n"]
    n_checklists = c.execute("SELECT COUNT(*) FROM (SELECT 1 FROM checklist_items "
                             "GROUP BY product_id, avl_id)").fetchone()[0]
    n_combos = c.execute("SELECT (SELECT COUNT(*) FROM avls WHERE active=1) * "
                         "(SELECT COUNT(*) FROM products WHERE active=1)").fetchone()[0]
    c.close()
    return templates.TemplateResponse(request, "files.html", {"user": user, "products": products,
        "avls": avls, "calls": calls_, "checks": checks, "atts": atts, "names": names,
        "picker": picker, "f_kind": kind, "f_ref": ref, "bundle_kinds": BUNDLE_KINDS,
        "ie_rows": ie_rows, "ie_picker": ie_picker, "ie_bundle_kinds": IE_BUNDLE_KINDS,
        "counts": counts, "n_checklists": n_checklists, "n_combos": n_combos})

@app.post("/files/upload")
async def upload(request: Request, kind: str = Form(...), ref_id: int = Form(...),
                 f: UploadFile = File(...), next_url: str = Form(""),
                 user=Depends(require_editor)):
    dest = next_url if next_url.startswith("/") else "/files"
    if kind not in ATT_KINDS:
        return RedirectResponse(dest, status_code=303)
    safe = _re.sub(r"[^A-Za-z0-9._-]", "_", f.filename or "file")
    stamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    stored = os.path.join(UPLOAD_DIR, f"{stamp}_{safe}")
    # Streamed in chunks: a 512 MB instance cannot afford to hold a large DNV
    # report or test-report PDF in memory just to write it to disk.
    with open(stored, "wb") as out:
        while chunk := await f.read(1024 * 1024):
            out.write(chunk)
    c = db.conn()
    c.execute("INSERT INTO attachments(kind, ref_id, filename, stored_path, uploaded_by, uploaded_at) "
              "VALUES(?,?,?,?,?,?)", (kind, ref_id, safe, stored, user["email"], now()))
    c.commit(); c.close()
    db.log(user["email"], "file:upload", f"{kind}#{ref_id} {safe}")
    return RedirectResponse(dest, status_code=303)

@app.post("/files/{att_id}/reassign")
def reassign(att_id: int, request: Request, target: str = Form(...),
             next_url: str = Form(""), user=Depends(require_editor)):
    """Re-file an attachment onto a different AVL / product / workstream requirement.

    `target` is "<kind>:<id>" so the kind and the row always move together.
    """
    dest = next_url if next_url.startswith("/") else "/files"
    kind, _, rid = target.partition(":")
    if kind not in ATT_KINDS or not rid.isdigit():
        return RedirectResponse(dest, status_code=303)
    ref_id = int(rid)
    c = db.conn()
    row = c.execute("SELECT * FROM attachments WHERE id=?", (att_id,)).fetchone()
    if not row:
        c.close()
        return RedirectResponse(dest, status_code=303)
    table = {"product": "products", "avl": "avls", "call": "calls", "checklist": "checklist_items"}[kind]
    if not c.execute(f"SELECT 1 FROM {table} WHERE id=?", (ref_id,)).fetchone():
        c.close()
        return RedirectResponse(dest + ("&" if "?" in dest else "?") + "err=badref", status_code=303)
    c.execute("UPDATE attachments SET kind=?, ref_id=? WHERE id=?", (kind, ref_id, att_id))
    c.commit(); c.close()
    db.log(user["email"], "file:reassign",
           f"{row['filename']}: {row['kind']}#{row['ref_id']} -> {kind}#{ref_id}")
    return RedirectResponse(dest, status_code=303)

@app.post("/files/{att_id}/delete")
def delete_file(att_id: int, request: Request, next_url: str = Form(""),
                user=Depends(require_editor)):
    dest = next_url if next_url.startswith("/") else "/files"
    c = db.conn()
    row = c.execute("SELECT * FROM attachments WHERE id=?", (att_id,)).fetchone()
    if row:
        c.execute("DELETE FROM attachments WHERE id=?", (att_id,))
        c.commit()
        try:
            os.remove(row["stored_path"])
        except OSError:
            pass
        db.log(user["email"], "file:delete", row["filename"])
    c.close()
    return RedirectResponse(dest, status_code=303)

BUNDLE_KINDS = {
    "checklist": "One workstream requirement",
    "product": "Product (files filed at product level)",
    "avl": "TPO AVL (files filed at AVL level)",
    "call": "Call",
    "product_all": "Everything for a product (all AVLs)",
    "avl_all": "Everything for a TPO AVL (all products + calls)",
}
IE_BUNDLE_KINDS = {
    "ie_item": "One IE review item",
    "ie_section": "One IE report section",
    "ie_report": "Whole IE report (evidence pack)",
}

def _bundle_files(c, kind, ref_id):
    """(label, [(folder, attachment_row)]) for a download bundle."""
    def plain(k, table, col="name"):
        row = c.execute(f"SELECT {col} AS n FROM {table} WHERE id=?", (ref_id,)).fetchone()
        if not row:
            return None, []
        atts = c.execute("SELECT * FROM attachments WHERE kind=? AND ref_id=? ORDER BY id",
                         (k, ref_id)).fetchall()
        return row["n"], [("", a) for a in atts]

    if kind == "product":
        return plain("product", "products")
    if kind == "avl":
        return plain("avl", "avls")
    if kind == "call":
        row = c.execute("SELECT calls.call_date, a.name AS avl_name FROM calls "
                        "JOIN avls a ON a.id=calls.avl_id WHERE calls.id=?", (ref_id,)).fetchone()
        if not row:
            return None, []
        atts = c.execute("SELECT * FROM attachments WHERE kind='call' AND ref_id=? ORDER BY id",
                         (ref_id,)).fetchall()
        return f"{row['avl_name']} {row['call_date']}", [("", a) for a in atts]
    if kind == "checklist":
        row = c.execute("SELECT ci.*, p.name AS product, a.name AS avl_name FROM checklist_items ci "
                        "JOIN products p ON p.id=ci.product_id JOIN avls a ON a.id=ci.avl_id "
                        "WHERE ci.id=?", (ref_id,)).fetchone()
        if not row:
            return None, []
        atts = c.execute("SELECT * FROM attachments WHERE kind='checklist' AND ref_id=? ORDER BY id",
                         (ref_id,)).fetchall()
        return f"{row['avl_name']} {row['product']} - {row['workstream']}", [("", a) for a in atts]

    if kind == "ie_item":
        row = c.execute("SELECT i.item_id, i.sub_section, s.title AS section, p.name AS product "
                        "FROM ie_report_items i JOIN ie_report_sections s ON s.id=i.section_id "
                        "JOIN ie_reports r ON r.id=i.report_id JOIN products p ON p.id=r.product_id "
                        "WHERE i.id=?", (ref_id,)).fetchone()
        if not row:
            return None, []
        atts = c.execute("SELECT * FROM attachments WHERE kind='ie_item' AND ref_id=? ORDER BY id",
                         (ref_id,)).fetchall()
        return f"{row['product']} {row['section']} {row['item_id']}", [("", a) for a in atts]

    out = []
    if kind in ("ie_section", "ie_report"):
        if kind == "ie_section":
            head = c.execute("SELECT s.title AS n, p.name AS product FROM ie_report_sections s "
                             "JOIN ie_reports r ON r.id=s.report_id JOIN products p ON p.id=r.product_id "
                             "WHERE s.id=?", (ref_id,)).fetchone()
            where, args = "s.id=?", (ref_id,)
        else:
            head = c.execute("SELECT r.name AS n, p.name AS product FROM ie_reports r "
                             "JOIN products p ON p.id=r.product_id WHERE r.id=?", (ref_id,)).fetchone()
            where, args = "r.id=?", (ref_id,)
        if not head:
            return None, []
        for r in c.execute(
                "SELECT att.*, s.title AS section, s.sort_order AS s_ord, i.item_id, i.sub_section "
                "FROM attachments att JOIN ie_report_items i ON i.id=att.ref_id "
                "JOIN ie_report_sections s ON s.id=i.section_id "
                "JOIN ie_reports r ON r.id=i.report_id "
                f"WHERE att.kind='ie_item' AND {where} ORDER BY s.sort_order, i.sort_order", args):
            folder = (f"{r['s_ord']+1:02d}_{r['section']}/{r['item_id']} {r['sub_section']}".strip()
                      if kind == "ie_report" else f"{r['item_id']} {r['sub_section']}".strip())
            out.append((folder, r))
        return head["n"], out

    # Roll-ups: everything filed anywhere under one product or one AVL.
    if kind == "product_all":
        row = c.execute("SELECT name FROM products WHERE id=?", (ref_id,)).fetchone()
        if not row:
            return None, []
        for a in c.execute("SELECT * FROM attachments WHERE kind='product' AND ref_id=? ORDER BY id",
                           (ref_id,)):
            out.append(("00_product-level", a))
        for r in c.execute(
                "SELECT att.*, av.name AS avl_name, ci.doc_category, ci.workstream "
                "FROM attachments att JOIN checklist_items ci ON ci.id=att.ref_id "
                "JOIN avls av ON av.id=ci.avl_id "
                "WHERE att.kind='checklist' AND ci.product_id=? ORDER BY av.name, ci.sort_order",
                (ref_id,)):
            out.append((f"{r['avl_name']}/{r['doc_category'] or 'Other'}", r))
        return row["name"], out

    if kind == "avl_all":
        row = c.execute("SELECT name FROM avls WHERE id=?", (ref_id,)).fetchone()
        if not row:
            return None, []
        for a in c.execute("SELECT * FROM attachments WHERE kind='avl' AND ref_id=? ORDER BY id",
                           (ref_id,)):
            out.append(("00_avl-level", a))
        for r in c.execute(
                "SELECT att.*, cl.call_date FROM attachments att JOIN calls cl ON cl.id=att.ref_id "
                "WHERE att.kind='call' AND cl.avl_id=? ORDER BY cl.call_date", (ref_id,)):
            out.append((f"01_calls/{r['call_date']}", r))
        for r in c.execute(
                "SELECT att.*, p.name AS product, ci.doc_category FROM attachments att "
                "JOIN checklist_items ci ON ci.id=att.ref_id JOIN products p ON p.id=ci.product_id "
                "WHERE att.kind='checklist' AND ci.avl_id=? ORDER BY p.name, ci.sort_order",
                (ref_id,)):
            out.append((f"{r['product']}/{r['doc_category'] or 'Other'}", r))
        return row["name"], out
    return None, []

@app.get("/files/bundle.zip")
def files_bundle(request: Request, kind: str = "", ref_id: int = 0, user=Depends(require_user)):
    """Download side of the uploader: the same targets, zipped with a manifest."""
    if (kind not in BUNDLE_KINDS and kind not in IE_BUNDLE_KINDS) or not ref_id:
        return RedirectResponse("/files?err=badtarget", status_code=303)
    c = db.conn()
    label, items = _bundle_files(c, kind, ref_id)
    c.close()
    if label is None:
        return RedirectResponse("/files?err=badtarget", status_code=303)
    if not items:
        return RedirectResponse(f"/files?err=empty&kind={kind}&ref={ref_id}", status_code=303)

    stamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    base = f"{_safe(label, 60)}_{stamp}"
    path = os.path.join(PKG_DIR, base + ".zip")
    rows, missing = [], 0
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        seen = set()
        for folder, a in items:
            if not os.path.exists(a["stored_path"]):
                missing += 1
                rows.append([folder or "/", a["filename"], a["uploaded_by"],
                             (a["uploaded_at"] or "")[:16], "FILE NOT FOUND ON DISK"])
                continue
            arc = f"{_safe_path(folder)}/{_safe(a['filename'], 90)}" if folder else _safe(a["filename"], 90)
            n = 1
            while arc in seen:            # same filename twice under one folder
                stem, _, ext = arc.rpartition(".")
                arc = f"{stem}_{n}.{ext}" if stem else f"{arc}_{n}"
                n += 1
            seen.add(arc)
            z.write(a["stored_path"], arc)
            rows.append([folder or "/", a["filename"], a["uploaded_by"],
                         (a["uploaded_at"] or "")[:16], arc])
        buf = io.StringIO()
        w = csv.writer(buf)
        w.writerow(["Folder", "File", "Uploaded by", "Uploaded at", "Path in zip"])
        for r in rows:
            w.writerow(r)
        z.writestr("MANIFEST.csv", buf.getvalue())

    kind_label = {**BUNDLE_KINDS, **IE_BUNDLE_KINDS}[kind]
    db.log(user["email"], "files:bundle",
           f"{kind_label} '{label}' - {len(items) - missing} file(s)")
    return FileResponse(path, filename=base + ".zip", media_type="application/zip")

@app.get("/files/{att_id}/download")
def download(att_id: int, request: Request, user=Depends(require_user)):
    c = db.conn()
    row = c.execute("SELECT * FROM attachments WHERE id=?", (att_id,)).fetchone()
    c.close()
    if not row:
        return RedirectResponse("/files", status_code=303)
    return FileResponse(row["stored_path"], filename=row["filename"])

# ---------------- 10) TPO-side contacts ----------------
@app.get("/contacts", response_class=HTMLResponse)
def contacts(request: Request, avl: int = 0, show_inactive: int = 0, q: str = "",
             user=Depends(require_user)):
    c = db.conn()
    avls = c.execute("SELECT id, name, account_manager FROM avls WHERE active=1 ORDER BY name").fetchall()
    sql = ("SELECT ct.*, a.name AS avl_name FROM contacts ct JOIN avls a ON a.id=ct.avl_id WHERE 1=1 ")
    params = []
    if avl:
        sql += "AND ct.avl_id=? "
        params.append(avl)
    if not show_inactive:
        sql += "AND ct.active=1 "
    if q.strip():
        sql += "AND (ct.name LIKE ? OR ct.role LIKE ? OR ct.email LIKE ?) "
        params += [f"%{q.strip()}%"] * 3
    rows = c.execute(sql + "ORDER BY a.name, ct.active DESC, ct.is_primary DESC, ct.name",
                     params).fetchall()
    # Group by AVL so the page reads as an address book, and show the empty ones.
    groups, seen = [], {}
    for a in avls:
        if avl and a["id"] != avl:
            continue
        seen[a["id"]] = {"avl": a, "rows": []}
        groups.append(seen[a["id"]])
    for r in rows:
        if r["avl_id"] in seen:
            seen[r["avl_id"]]["rows"].append(r)
    n_live = sum(1 for r in rows if r["active"])
    c.close()
    return templates.TemplateResponse(request, "contacts.html", {"user": user, "avls": avls,
        "groups": groups, "sel_a": avl, "q": q, "show_inactive": show_inactive,
        "contact_roles": db.CONTACT_ROLES, "n_live": n_live})

@app.post("/contacts/add")
def contact_add(request: Request, avl_id: int = Form(...), name: str = Form(...),
                role: str = Form(""), email: str = Form(""), phone: str = Form(""),
                website: str = Form(""), notes: str = Form(""), is_primary: int = Form(0),
                user=Depends(require_editor)):
    nm = name.strip()
    if not nm:
        return RedirectResponse(f"/contacts?avl={avl_id}", status_code=303)
    c = db.conn()
    if is_primary:
        c.execute("UPDATE contacts SET is_primary=0 WHERE avl_id=?", (avl_id,))
    c.execute("INSERT INTO contacts(avl_id, name, role, email, phone, website, notes, is_primary, "
              "created_by, created_at, updated_by, updated_at) VALUES(?,?,?,?,?,?,?,?,?,?,?,?)",
              (avl_id, nm, role.strip(), email.strip().lower(), phone.strip(),
               db.normalize_website(website), notes.strip(), 1 if is_primary else 0,
               user["email"], now(), user["email"], now()))
    aname = c.execute("SELECT name FROM avls WHERE id=?", (avl_id,)).fetchone()["name"]
    c.commit(); c.close()
    db.log(user["email"], "contact:add", nm, avl_id=avl_id, entity="contact")
    return RedirectResponse(f"/contacts?avl={avl_id}", status_code=303)

@app.post("/contacts/{cid}/save")
def contact_save(cid: int, request: Request, name: str = Form(...), role: str = Form(""),
                 email: str = Form(""), phone: str = Form(""), website: str = Form(""),
                 notes: str = Form(""), user=Depends(require_editor)):
    nm = name.strip()
    c = db.conn()
    row = c.execute("SELECT avl_id FROM contacts WHERE id=?", (cid,)).fetchone()
    if not row or not nm:
        c.close()
        return RedirectResponse("/contacts", status_code=303)
    c.execute("UPDATE contacts SET name=?, role=?, email=?, phone=?, website=?, notes=?, "
              "updated_by=?, updated_at=? WHERE id=?",
              (nm, role.strip(), email.strip().lower(), phone.strip(),
               db.normalize_website(website), notes.strip(), user["email"], now(), cid))
    c.commit(); c.close()
    db.log(user["email"], "contact:save", nm)
    return RedirectResponse(f"/contacts?avl={row['avl_id']}", status_code=303)

@app.post("/contacts/{cid}/primary")
def contact_primary(cid: int, request: Request, user=Depends(require_editor)):
    """Exactly one primary point of contact per AVL."""
    c = db.conn()
    row = c.execute("SELECT avl_id, name, is_primary FROM contacts WHERE id=?", (cid,)).fetchone()
    if not row:
        c.close()
        return RedirectResponse("/contacts", status_code=303)
    c.execute("UPDATE contacts SET is_primary=0 WHERE avl_id=?", (row["avl_id"],))
    if not row["is_primary"]:
        c.execute("UPDATE contacts SET is_primary=1 WHERE id=?", (cid,))
    c.commit(); c.close()
    db.log(user["email"], "contact:primary",
           f"{row['name']}" + (" (cleared)" if row["is_primary"] else ""))
    return RedirectResponse(f"/contacts?avl={row['avl_id']}", status_code=303)

@app.post("/contacts/{cid}/toggle")
def contact_toggle(cid: int, request: Request, user=Depends(require_editor)):
    """Soft delete: people leave TPOs, but the call history should still name them."""
    c = db.conn()
    row = c.execute("SELECT avl_id, name, active FROM contacts WHERE id=?", (cid,)).fetchone()
    if not row:
        c.close()
        return RedirectResponse("/contacts", status_code=303)
    c.execute("UPDATE contacts SET active = 1 - active, is_primary = CASE WHEN active=1 THEN 0 "
              "ELSE is_primary END, updated_by=?, updated_at=? WHERE id=?",
              (user["email"], now(), cid))
    c.commit(); c.close()
    db.log(user["email"], "contact:toggle",
           f"{row['name']} -> {'inactive' if row['active'] else 'active'}")
    return RedirectResponse(f"/contacts?avl={row['avl_id']}", status_code=303)

@app.get("/contacts.csv")
def contacts_csv(request: Request, user=Depends(require_user)):
    c = db.conn()
    rows = c.execute("SELECT a.name AS avl_name, ct.name, ct.role, ct.email, ct.phone, ct.website, "
                     "ct.is_primary, ct.notes FROM contacts ct JOIN avls a ON a.id=ct.avl_id "
                     "WHERE ct.active=1 ORDER BY a.name, ct.is_primary DESC, ct.name").fetchall()
    c.close()
    return _csv_response([[r["avl_name"], r["name"], r["role"], r["email"], r["phone"],
                           r["website"], "Yes" if r["is_primary"] else "", r["notes"]] for r in rows],
                         ["TPO AVL", "Name", "Role", "Email", "Phone", "Website", "Primary", "Notes"],
                         "avl_contacts.csv")

# ---------------- 11) dataroom package builder ----------------
import zipfile, json

PKG_DIR = os.path.join(UPLOAD_DIR, "packages")
os.makedirs(PKG_DIR, exist_ok=True)

def _safe(name, limit=80):
    """One path segment, safe on Windows and macOS. & and , are legal, / is not."""
    s = _re.sub(r"[^A-Za-z0-9 ._()&,+-]", "_", (name or "").strip())
    return (_re.sub(r"_+", "_", s)[:limit] or "item").strip(" ._")

def _safe_path(path, limit=60):
    """Sanitise each segment so 'Product/Category' nests instead of flattening."""
    return "/".join(_safe(seg, limit) for seg in str(path or "").split("/") if seg.strip())

def _package_contents(c, product_id, avl_id, scope):
    """Every in-scope requirement with its files, in checklist order.

    Requirements with no file are kept and marked as gaps - the content list is
    only useful if it shows what is still missing.
    """
    frag, extra = db.scope_filter(scope)
    rows = c.execute("SELECT * FROM checklist_items WHERE product_id=? AND avl_id=? " + frag +
                     "ORDER BY sort_order, id", [product_id, avl_id] + extra).fetchall()
    atts = _checklist_atts(c, [r["id"] for r in rows])
    ie_prog = db.ie_report_progress(c, [r["ie_report_id"] for r in rows])
    groups, seen = [], {}
    for r in rows:
        key = r["doc_category"] or "Other"
        if key not in seen:
            seen[key] = {"name": key, "items": []}
            groups.append(seen[key])
        files = atts.get(r["id"], [])
        seen[key]["items"].append({"row": r, "files": files, "gap": not files,
                                   "ie": ie_prog.get(r["ie_report_id"])})
    return groups

def _dominant_template(c, product_id, avl_id):
    row = c.execute("SELECT template_id, COUNT(*) n FROM checklist_items "
                    "WHERE product_id=? AND avl_id=? AND template_id IS NOT NULL "
                    "GROUP BY template_id ORDER BY n DESC LIMIT 1", (product_id, avl_id)).fetchone()
    return row["template_id"] if row else None

def _template_drift(c, product_id, avl_id, template_id):
    """Template requirements that never made it onto this checklist.

    Seeding is a point-in-time copy, so a template edited afterwards leaves the
    checklist behind. Those are gaps too, and worth naming in a submission.
    """
    if not template_id:
        return []
    have = {r["workstream"] for r in c.execute(
        "SELECT workstream FROM checklist_items WHERE product_id=? AND avl_id=?",
        (product_id, avl_id))}
    return [r for r in c.execute(
        "SELECT doc_category, workstream, obligation, description "
        "FROM workstream_template_items WHERE template_id=? ORDER BY sort_order, id",
        (template_id,)) if r["workstream"] not in have]

def _package_stats(groups):
    items = [i for g in groups for i in g["items"]]
    n_gaps = sum(1 for i in items if i["gap"])
    # A gap covered by a linked IE review is being worked somewhere the reader can
    # check, which is a different thing from nobody having started.
    covered = sum(1 for i in items if i["gap"] and i.get("ie"))
    return {"n_reqs": len(items),
            "n_files": sum(len(i["files"]) for i in items),
            "n_gaps": n_gaps, "ie_covered": covered,
            "blocking": sum(1 for i in items
                            if i["gap"] and not i.get("ie") and i["row"]["obligation"] == "Required")}

@app.get("/dataroom/package", response_class=HTMLResponse)
def package_preview(request: Request, product: int = 0, avl: int = 0, scope: str = "all",
                    template: int = 0, user=Depends(require_user)):
    if scope not in db.PACKAGE_SCOPES:
        scope = "all"
    c = db.conn()
    products = c.execute("SELECT id, name, category FROM products WHERE active=1 "
                         "ORDER BY category, name").fetchall()
    avls = c.execute("SELECT id, name FROM avls WHERE active=1 ORDER BY name").fetchall()
    product = product or (products[0]["id"] if products else 0)
    avl = avl or (avls[0]["id"] if avls else 0)
    groups = _package_contents(c, product, avl, scope)
    stats = _package_stats(groups)
    tmpl_id = template or _dominant_template(c, product, avl)
    drift = _template_drift(c, product, avl, tmpl_id)
    tmpl = c.execute("SELECT id, name FROM workstream_templates WHERE id=?", (tmpl_id,)).fetchone() \
           if tmpl_id else None
    all_tmpls = db.templates_for(c, "", avl) or c.execute(
        "SELECT id, name, category FROM workstream_templates WHERE active=1 ORDER BY name").fetchall()
    prior = c.execute("SELECT * FROM packages WHERE product_id=? AND avl_id=? "
                      "ORDER BY id DESC LIMIT 20", (product, avl)).fetchall()
    next_rev = (c.execute("SELECT COALESCE(MAX(revision), 0) + 1 AS n FROM packages "
                          "WHERE product_id=? AND avl_id=?", (product, avl)).fetchone()["n"])
    pname = c.execute("SELECT name FROM products WHERE id=?", (product,)).fetchone()
    aname = c.execute("SELECT name FROM avls WHERE id=?", (avl,)).fetchone()
    c.close()
    return templates.TemplateResponse(request, "package.html", {"user": user, "products": products,
        "avls": avls, "sel_p": product, "sel_a": avl, "scope": scope, "scopes": db.PACKAGE_SCOPES,
        "groups": groups, "stats": stats, "prior": prior, "drift": drift, "tmpl": tmpl,
        "all_tmpls": all_tmpls, "next_rev": next_rev, "today": datetime.date.today().isoformat(),
        "product_name": pname["name"] if pname else "", "avl_name": aname["name"] if aname else ""})

@app.post("/dataroom/package/build")
def package_build(request: Request, product_id: int = Form(...), avl_id: int = Form(...),
                  scope: str = Form("all"), label: str = Form(""), template_id: str = Form(""),
                  rev_date: str = Form(""), user=Depends(require_editor)):
    """Zip the attached files with a manifest and a dataroom summary.

    The summary compares the checklist against the template it was seeded from,
    so it names three things: what is included, what is on the checklist with no
    document, and what the template asks for that the checklist never picked up.
    """
    if scope not in db.PACKAGE_SCOPES:
        scope = "all"
    c = db.conn()
    groups = _package_contents(c, product_id, avl_id, scope)
    stats = _package_stats(groups)
    tid = int(template_id) if str(template_id).strip().isdigit() and int(template_id) > 0 \
          else _dominant_template(c, product_id, avl_id)
    drift = _template_drift(c, product_id, avl_id, tid)
    trow = c.execute("SELECT name FROM workstream_templates WHERE id=?", (tid,)).fetchone() if tid else None
    tname = trow["name"] if trow else "(no template recorded)"
    pname = c.execute("SELECT name FROM products WHERE id=?", (product_id,)).fetchone()["name"]
    aname = c.execute("SELECT name FROM avls WHERE id=?", (avl_id,)).fetchone()["name"]

    rev = c.execute("SELECT COALESCE(MAX(revision), 0) + 1 AS n FROM packages "
                    "WHERE product_id=? AND avl_id=?", (product_id, avl_id)).fetchone()["n"]
    rdate = rev_date.strip() or datetime.date.today().isoformat()
    rev_code = f"R{rev:02d}"
    stamp = rdate.replace("-", "")
    base = f"{_safe(aname, 30)}_{_safe(pname, 30)}_{rev_code}_{stamp}"
    stored = os.path.join(PKG_DIR, base + ".zip")
    full_code = f"{aname} / {pname} / {rev_code} / {rdate}"

    manifest_rows, summary_rows, missing_paths = [], [], []
    with zipfile.ZipFile(stored, "w", zipfile.ZIP_DEFLATED) as z:
        for gi, g in enumerate(groups, 1):
            folder = f"{gi:02d}_{_safe(g['name'], 50)}"
            for item in g["items"]:
                r = item["row"]
                if item["gap"]:
                    ie = item.get("ie")
                    res = (f"IE REVIEW - {ie['reviewer']} {ie['status']}, "
                           f"{ie['done']}/{ie['n']} accepted") if ie else "NO DOCUMENT"
                    summary_rows.append([g["name"], r["workstream"], r["obligation"], r["status"],
                                         0, res])
                    manifest_rows.append([g["name"], r["workstream"], r["obligation"],
                                          r["status"], "",
                                          f"COVERED BY IE REVIEW ({ie['name']})" if ie
                                          else "MISSING - no file attached"])
                    continue
                sub = f"{folder}/{_safe(r['workstream'], 60)}"
                got = 0
                for f in item["files"]:
                    arc = f"{sub}/{_safe(f['filename'], 90)}"
                    if os.path.exists(f["stored_path"]):
                        z.write(f["stored_path"], arc)
                        got += 1
                        manifest_rows.append([g["name"], r["workstream"], r["obligation"],
                                              r["status"], arc, ""])
                    else:
                        missing_paths.append(f["filename"])
                        manifest_rows.append([g["name"], r["workstream"], r["obligation"],
                                              r["status"], "", "FILE NOT FOUND ON DISK"])
                summary_rows.append([g["name"], r["workstream"], r["obligation"], r["status"],
                                     got, "INCLUDED" if got else "NO DOCUMENT"])
        for d in drift:
            summary_rows.append([d["doc_category"], d["workstream"], d["obligation"],
                                 "-", 0, "NOT ON CHECKLIST"])

        def csv_bytes(header, rows):
            buf = io.StringIO()
            w = csv.writer(buf)
            w.writerow(header)
            for row in rows:
                w.writerow(row)
            return buf.getvalue()

        z.writestr("MANIFEST.csv", csv_bytes(
            ["Document Category", "Required Document", "Obligation", "Status",
             "File in package", "Gap"], manifest_rows))
        z.writestr("SUMMARY.csv", csv_bytes(
            ["Document Category", "Required Document", "Obligation", "Checklist status",
             "Files included", "Result"], summary_rows))

        # readable summary
        by_cat = {}
        for cat, req, obl, st_, n, res in summary_rows:
            d = by_cat.setdefault(cat or "Other", {"n": 0, "inc": 0, "nodoc": 0, "untracked": 0, "files": 0})
            d["n"] += 1
            d["files"] += n
            d["inc"] += res == "INCLUDED"
            d["nodoc"] += res == "NO DOCUMENT"
            d["untracked"] += res == "NOT ON CHECKLIST"
        L = [f"DATAROOM SUBMISSION SUMMARY",
             f"{'=' * 60}",
             f"TPO AVL          : {aname}",
             f"Product          : {pname}",
             f"Revision         : {rev_code}",
             f"Revision date    : {rdate}",
             f"Built            : {now()} by {user['email']}",
             f"Scope            : {db.PACKAGE_SCOPES[scope]}",
             f"Compared against : {tname}"]
        if label.strip():
            L.append(f"Label            : {label.strip()}")
        L += ["",
              f"{stats['n_files']} document(s) included against {stats['n_reqs']} requirement(s) in scope.",
              f"{stats['n_gaps']} requirement(s) have no document "
              f"({stats['ie_covered']} of those covered by a linked IE review, "
              f"{stats['blocking']} Required and uncovered).",
              f"{len(drift)} template requirement(s) are not on this checklist yet.", "",
              "BY DOCUMENT CATEGORY", "-" * 60,
              f"{'Category':<38}{'Reqs':>5}{'Incl':>6}{'NoDoc':>7}{'NotTrk':>7}"]
        for cat, d in by_cat.items():
            L.append(f"{cat[:37]:<38}{d['n']:>5}{d['inc']:>6}{d['nodoc']:>7}{d['untracked']:>7}")
        tot = {k: sum(d[k] for d in by_cat.values()) for k in ("n", "inc", "nodoc", "untracked")}
        L += [f"{'TOTAL':<38}{tot['n']:>5}{tot['inc']:>6}{tot['nodoc']:>7}{tot['untracked']:>7}", ""]

        L += ["CONTENTS", "-" * 60]
        for gi, g in enumerate(groups, 1):
            got = sum(len(i["files"]) for i in g["items"])
            L.append(f"{gi:02d}. {g['name']}  -  {got} file(s)")
            for item in g["items"]:
                r = item["row"]
                ie = item.get("ie")
                mark = ("  [~] IE REVIEW" if ie else "  [ ] NO DOCUMENT") if item["gap"] else "  [x]"
                L.append(f"    {mark} ({r['obligation']}, {r['status']}) {r['workstream']}")
                if ie:
                    L.append(f"          covered by: {ie['name']} - {ie['status']}, "
                             f"{ie['done']}/{ie['n']} items accepted")
                for f in item["files"]:
                    L.append(f"          - {f['filename']}")
            L.append("")
        if drift:
            L += ["NOT ON THIS CHECKLIST", "-" * 60,
                  f"In '{tname}' but never seeded here. Re-apply the template on the Dataroom",
                  "page (Merge) to start tracking them.", ""]
            for d in drift:
                L.append(f"    [ ] ({d['obligation']}) {d['doc_category']}: {d['workstream']}")
            L.append("")
        if missing_paths:
            L += ["WARNING - recorded attachments whose file was not on disk:",
                  *[f"  - {m}" for m in missing_paths], ""]
        z.writestr("DATAROOM_SUMMARY.txt", "\n".join(L))

    size = os.path.getsize(stored)
    c.execute("INSERT INTO packages(product_id, avl_id, label, scope, n_files, n_reqs, n_gaps, "
              "n_untracked, bytes, stored_path, manifest, revision, rev_code, rev_date, "
              "template_id, created_by, created_at) VALUES(?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)",
              (product_id, avl_id, label.strip(), scope, stats["n_files"], stats["n_reqs"],
               stats["n_gaps"], len(drift), size, stored, json.dumps(manifest_rows), rev,
               rev_code, rdate, tid, user["email"], now()))
    c.execute("UPDATE packages SET superseded=1 WHERE product_id=? AND avl_id=? AND revision<?",
              (product_id, avl_id, rev))
    c.commit(); c.close()
    db.log(user["email"], "package:build",
           f"{rev_code} ({rdate}) - {stats['n_files']} files, {stats['n_gaps']} gaps, "
           f"{len(drift)} untracked", avl_id=avl_id, product_id=product_id, entity="package")
    return FileResponse(stored, filename=base + ".zip", media_type="application/zip")

@app.get("/dataroom/package/{pkg_id}/download")
def package_download(pkg_id: int, request: Request, user=Depends(require_user)):
    c = db.conn()
    row = c.execute("SELECT * FROM packages WHERE id=?", (pkg_id,)).fetchone()
    c.close()
    if not row or not os.path.exists(row["stored_path"]):
        return RedirectResponse("/dataroom/package?err=gone", status_code=303)
    return FileResponse(row["stored_path"], filename=os.path.basename(row["stored_path"]),
                        media_type="application/zip")

@app.get("/dataroom/package/{pkg_id}/manifest.csv")
def package_manifest(pkg_id: int, request: Request, user=Depends(require_user)):
    """The content list as sent, even if files have changed since."""
    c = db.conn()
    row = c.execute("SELECT * FROM packages WHERE id=?", (pkg_id,)).fetchone()
    c.close()
    if not row:
        return RedirectResponse("/dataroom/package", status_code=303)
    rows = json.loads(row["manifest"] or "[]")
    return _csv_response(rows, ["Document Category", "Required Document", "Obligation",
                                "Status", "File in package", "Gap"],
                         f"manifest_{pkg_id}.csv")

@app.post("/dataroom/package/{pkg_id}/delete")
def package_delete(pkg_id: int, request: Request, user=Depends(require_editor)):
    c = db.conn()
    row = c.execute("SELECT * FROM packages WHERE id=?", (pkg_id,)).fetchone()
    if row:
        c.execute("DELETE FROM packages WHERE id=?", (pkg_id,))
        c.commit()
        try:
            os.remove(row["stored_path"])
        except OSError:
            pass
        db.log(user["email"], "package:delete", os.path.basename(row["stored_path"]))
    c.close()
    return RedirectResponse(f"/dataroom/package?product={row['product_id']}&avl={row['avl_id']}"
                            if row else "/dataroom/package", status_code=303)

# ---------------- 12) IE / DNV technology reviews ----------------
# Kept separate from the AVL workstream templates on purpose: an IE review is a
# document a third party writes about a product, not a per-TPO dataroom checklist.
IE_KIND = "ie_item"
IE_SEC_KIND = "ie_section"   # evidence filed against a whole section

def _ie_progress(c, report_id):
    """Per-section rollup plus totals for one report."""
    secs = c.execute("SELECT * FROM ie_report_sections WHERE report_id=? ORDER BY sort_order, id",
                     (report_id,)).fetchall()
    items = c.execute("SELECT * FROM ie_report_items WHERE report_id=? ORDER BY sort_order, id",
                      (report_id,)).fetchall()
    sec_atts = {}
    if secs:
        qs = ",".join("?" * len(secs))
        for a in c.execute(f"SELECT * FROM attachments WHERE kind='{IE_SEC_KIND}' "
                           f"AND ref_id IN ({qs}) ORDER BY id DESC", [s["id"] for s in secs]):
            sec_atts.setdefault(a["ref_id"], []).append(a)
    atts = {}
    if items:
        qs = ",".join("?" * len(items))
        for a in c.execute(f"SELECT * FROM attachments WHERE kind='{IE_KIND}' AND ref_id IN ({qs}) "
                           "ORDER BY id DESC", [i["id"] for i in items]):
            atts.setdefault(a["ref_id"], []).append(a)
    groups, by_sec = [], {}
    for s in secs:
        by_sec[s["id"]] = {"sec": s, "items": [], "n": 0, "done": 0, "wip": 0,
                           "blocked": 0, "open": 0, "crit_open": 0,
                           "files": len(sec_atts.get(s["id"], []))}
        groups.append(by_sec[s["id"]])
    for i in items:
        g = by_sec.get(i["section_id"])
        if not g:
            continue
        g["items"].append(i)
        g["n"] += 1
        g["files"] += len(atts.get(i["id"], []))
        st = i["status"]
        if st in ("Accepted", "N/A"):
            g["done"] += 1
        elif st == "Blocked":
            g["blocked"] += 1
        elif st in ("In Progress", "Submitted"):
            g["wip"] += 1
        else:
            g["open"] += 1
        if i["priority"] == "Critical" and st not in ("Accepted", "N/A"):
            g["crit_open"] += 1
    for g in groups:
        g["pct"] = round(100 * g["done"] / g["n"]) if g["n"] else 0
    totals = {k: sum(g[k] for g in groups) for k in
              ("n", "done", "wip", "blocked", "open", "files", "crit_open")}
    totals["pct"] = round(100 * totals["done"] / totals["n"]) if totals["n"] else 0
    return groups, totals, atts, sec_atts

@app.get("/ie", response_class=HTMLResponse)
def ie_home(request: Request, user=Depends(require_user)):
    c = db.conn()
    reports = c.execute(
        "SELECT r.*, p.name AS product, p.category, "
        "(SELECT COUNT(*) FROM ie_report_items i WHERE i.report_id=r.id) AS n_items, "
        "(SELECT COUNT(*) FROM ie_report_items i WHERE i.report_id=r.id "
        " AND i.status IN ('Accepted','N/A')) AS n_done "
        "FROM ie_reports r JOIN products p ON p.id=r.product_id "
        "WHERE r.active=1 ORDER BY r.id DESC").fetchall()
    tmpls = c.execute("SELECT t.*, "
                      "(SELECT COUNT(*) FROM ie_template_sections s WHERE s.template_id=t.id) AS n_sec, "
                      "(SELECT COUNT(*) FROM ie_template_items i JOIN ie_template_sections s "
                      " ON s.id=i.section_id WHERE s.template_id=t.id) AS n_items "
                      "FROM ie_templates t WHERE t.active=1 ORDER BY t.name").fetchall()
    products = c.execute("SELECT id, name, category FROM products WHERE active=1 "
                         "ORDER BY category, name").fetchall()
    c.close()
    return templates.TemplateResponse(request, "ie.html", {"user": user, "reports": reports,
        "tmpls": tmpls, "products": products, "reviewers": db.IE_REVIEWERS,
        "statuses": db.IE_REPORT_STATUSES})

@app.post("/ie/reports/add")
def ie_report_add(request: Request, product_id: int = Form(...), template_id: int = Form(...),
                  name: str = Form(""), reviewer: str = Form("DNV"),
                  kickoff_date: str = Form(""), target_date: str = Form(""),
                  user=Depends(require_editor)):
    c = db.conn()
    t = c.execute("SELECT * FROM ie_templates WHERE id=?", (template_id,)).fetchone()
    p = c.execute("SELECT name FROM products WHERE id=?", (product_id,)).fetchone()
    if not t or not p:
        c.close()
        return RedirectResponse("/ie?err=bad", status_code=303)
    rname = name.strip() or f"{p['name']} - {reviewer} Technology Review"
    c.execute("INSERT INTO ie_reports(product_id, template_id, name, reviewer, kickoff_date, "
              "target_date, created_by, created_at) VALUES(?,?,?,?,?,?,?,?)",
              (product_id, template_id, rname, reviewer, kickoff_date, target_date,
               user["email"], now()))
    rid = c.execute("SELECT last_insert_rowid() AS id").fetchone()["id"]
    for s in c.execute("SELECT * FROM ie_template_sections WHERE template_id=? "
                       "ORDER BY sort_order, id", (template_id,)).fetchall():
        c.execute("INSERT INTO ie_report_sections(report_id, code, title, owner, sort_order) "
                  "VALUES(?,?,?,?,?)", (rid, s["code"], s["title"], s["owner"], s["sort_order"]))
        sid = c.execute("SELECT last_insert_rowid() AS id").fetchone()["id"]
        for i in c.execute("SELECT * FROM ie_template_items WHERE section_id=? "
                           "ORDER BY sort_order, id", (s["id"],)).fetchall():
            c.execute("INSERT INTO ie_report_items(report_id, section_id, item_id, sub_section, "
                      "review_item, evidence, priority, source, owner, sort_order, "
                      "updated_by, updated_at) VALUES(?,?,?,?,?,?,?,?,?,?,?,?)",
                      (rid, sid, i["item_id"], i["sub_section"], i["review_item"], i["evidence"],
                       i["priority"], i["source"], i["suggested_owner"], i["sort_order"],
                       user["email"], now()))
    c.commit(); c.close()
    db.log(user["email"], "ie:report:add", f"{rname} from '{t['name']}'")
    return RedirectResponse(f"/ie/report/{rid}", status_code=303)

@app.get("/ie/report/{rid}", response_class=HTMLResponse)
def ie_report(request: Request, rid: int, only: str = "", user=Depends(require_user)):
    c = db.conn()
    r = c.execute("SELECT r.*, p.name AS product, p.category FROM ie_reports r "
                  "JOIN products p ON p.id=r.product_id WHERE r.id=?", (rid,)).fetchone()
    if not r:
        c.close()
        return RedirectResponse("/ie", status_code=303)
    groups, totals, atts, sec_atts = _ie_progress(c, rid)
    if only:
        for g in groups:
            if only == "open":
                g["items"] = [i for i in g["items"] if i["status"] not in ("Accepted", "N/A")]
            elif only == "critical":
                g["items"] = [i for i in g["items"] if i["priority"] == "Critical"]
            elif only == "blocked":
                g["items"] = [i for i in g["items"] if i["status"] == "Blocked"]
            elif only == "nofile":
                g["items"] = [i for i in g["items"] if not atts.get(i["id"])]
    people = c.execute("SELECT id, name, org FROM people WHERE active=1 ORDER BY name").fetchall()
    avls = c.execute("SELECT id, name FROM avls WHERE active=1 ORDER BY name").fetchall()
    # The other direction of the dataroom link: which TPO requirements this answers.
    satisfies = c.execute(
        "SELECT ci.id, ci.workstream, ci.doc_category, ci.status, ci.avl_id, ci.product_id, "
        "a.name AS avl_name, p.name AS product FROM checklist_items ci "
        "JOIN avls a ON a.id=ci.avl_id JOIN products p ON p.id=ci.product_id "
        "WHERE ci.ie_report_id=? ORDER BY a.name", (rid,)).fetchall()
    c.close()
    return templates.TemplateResponse(request, "ie_report.html", {"user": user, "r": r,
        "satisfies": satisfies,
        "groups": groups, "totals": totals, "atts": atts, "sec_atts": sec_atts,
        "people": people, "avls": avls,
        "item_statuses": db.IE_ITEM_STATUSES, "priorities": db.IE_PRIORITIES,
        "report_statuses": db.IE_REPORT_STATUSES, "reviewers": db.IE_REVIEWERS, "only": only})

@app.post("/ie/report/{rid}/save")
def ie_report_save(rid: int, request: Request, name: str = Form(...), reviewer: str = Form("DNV"),
                   status: str = Form("Planning"), kickoff_date: str = Form(""),
                   target_date: str = Form(""), shared_with: str = Form(""),
                   notes: str = Form(""), user=Depends(require_editor)):
    c = db.conn()
    c.execute("UPDATE ie_reports SET name=?, reviewer=?, status=?, kickoff_date=?, target_date=?, "
              "shared_with=?, notes=? WHERE id=?",
              (name.strip(), reviewer, status, kickoff_date, target_date,
               shared_with.strip(), notes.strip(), rid))
    c.commit(); c.close()
    db.log(user["email"], "ie:report:save", f"{name.strip()} -> {status}")
    return RedirectResponse(f"/ie/report/{rid}", status_code=303)

@app.post("/ie/report/{rid}/delete")
def ie_report_delete(rid: int, request: Request, user=Depends(require_admin)):
    c = db.conn()
    row = c.execute("SELECT name FROM ie_reports WHERE id=?", (rid,)).fetchone()
    c.execute("UPDATE ie_reports SET active=0 WHERE id=?", (rid,))
    c.commit(); c.close()
    if row:
        db.log(user["email"], "ie:report:delete", row["name"])
    return RedirectResponse("/ie", status_code=303)

@app.post("/ie/item/{iid}/save")
def ie_item_save(iid: int, request: Request, status: str = Form(...),
                 owner_person_id: str = Form(""), owner_other: str = Form(""),
                 due_date: str = Form(""), priority: str = Form("Normal"),
                 gap: str = Form(""), notes: str = Form(""), user=Depends(require_editor)):
    if status not in db.IE_ITEM_STATUSES:
        status = "Not Started"
    if priority not in db.IE_PRIORITIES:
        priority = "Normal"
    c = db.conn()
    row = c.execute("SELECT report_id, review_item, owner FROM ie_report_items WHERE id=?",
                    (iid,)).fetchone()
    if not row:
        c.close()
        return RedirectResponse("/ie", status_code=303)
    oid = int(owner_person_id) if owner_person_id.strip().isdigit() and int(owner_person_id) > 0 else None
    oname = ""
    if oid:
        p = c.execute("SELECT name FROM people WHERE id=? AND active=1", (oid,)).fetchone()
        oname, oid = (p["name"], oid) if p else ("", None)
    owner_txt = oname or owner_other.strip() or (row["owner"] if not oid else "")
    c.execute("UPDATE ie_report_items SET status=?, owner=?, owner_person_id=?, due_date=?, "
              "priority=?, gap=?, notes=?, updated_by=?, updated_at=? WHERE id=?",
              (status, owner_txt, oid, due_date, priority, gap.strip(), notes.strip(),
               user["email"], now(), iid))
    c.commit(); c.close()
    db.log(user["email"], "ie:item", f"{row['review_item'][:60]} -> {status}")
    return RedirectResponse(f"/ie/report/{row['report_id']}", status_code=303)

@app.post("/ie/report/{rid}/section/{sid}/save")
def ie_section_save(rid: int, sid: int, request: Request, owner_person_id: str = Form(""),
                    owner_other: str = Form(""), narrative: str = Form(None),
                    only: str = Form(""), user=Depends(require_editor)):
    c = db.conn()
    oid = int(owner_person_id) if owner_person_id.strip().isdigit() and int(owner_person_id) > 0 else None
    oname = ""
    if oid:
        p = c.execute("SELECT name FROM people WHERE id=? AND active=1", (oid,)).fetchone()
        oname, oid = (p["name"], oid) if p else ("", None)
    # Blank owner inputs mean "leave it alone", not "clear it" - the narrative
    # shares this form, so saving prose must not wipe the section owner.
    if oname or owner_other.strip():
        c.execute("UPDATE ie_report_sections SET owner=?, owner_person_id=? "
                  "WHERE id=? AND report_id=?",
                  (oname or owner_other.strip(), oid, sid, rid))
    if narrative is not None:
        c.execute("UPDATE ie_report_sections SET narrative=? WHERE id=? AND report_id=?",
                  (narrative.strip(), sid, rid))
    c.commit(); c.close()
    q = f"?only={only}" if only else ""
    return RedirectResponse(f"/ie/report/{rid}{q}#s{sid}", status_code=303)

@app.post("/ie/report/{rid}/item/add")
def ie_item_add(rid: int, request: Request, section_id: int = Form(...),
                review_item: str = Form(...), sub_section: str = Form(""),
                item_id: str = Form(""), evidence: str = Form(""),
                priority: str = Form("Normal"), user=Depends(require_editor)):
    ri = review_item.strip()
    if priority not in db.IE_PRIORITIES:
        priority = "Normal"
    if ri:
        c = db.conn()
        nxt = c.execute("SELECT COALESCE(MAX(sort_order), 0) + 1 FROM ie_report_items "
                        "WHERE report_id=? AND section_id=?", (rid, section_id)).fetchone()[0]
        c.execute("INSERT INTO ie_report_items(report_id, section_id, item_id, sub_section, "
                  "review_item, evidence, priority, sort_order, updated_by, updated_at) "
                  "VALUES(?,?,?,?,?,?,?,?,?,?)",
                  (rid, section_id, item_id.strip(), sub_section.strip(), ri, evidence.strip(),
                   priority, nxt, user["email"], now()))
        c.commit(); c.close()
        db.log(user["email"], "ie:item:add", ri[:60])
    return RedirectResponse(f"/ie/report/{rid}", status_code=303)

@app.post("/ie/item/{iid}/delete")
def ie_item_delete(iid: int, request: Request, user=Depends(require_editor)):
    c = db.conn()
    row = c.execute("SELECT report_id, review_item FROM ie_report_items WHERE id=?", (iid,)).fetchone()
    if not row:
        c.close()
        return RedirectResponse("/ie", status_code=303)
    n = c.execute(f"SELECT COUNT(*) FROM attachments WHERE kind='{IE_KIND}' AND ref_id=?",
                  (iid,)).fetchone()[0]
    if n:
        c.close()
        return RedirectResponse(f"/ie/report/{row['report_id']}?err=files", status_code=303)
    c.execute("DELETE FROM ie_report_items WHERE id=?", (iid,))
    c.commit(); c.close()
    db.log(user["email"], "ie:item:delete", row["review_item"][:60])
    return RedirectResponse(f"/ie/report/{row['report_id']}", status_code=303)

# ---------------- 12b) IE template management ----------------
@app.get("/ie/templates", response_class=HTMLResponse)
def ie_templates(request: Request, t: int = 0, user=Depends(require_user)):
    c = db.conn()
    tmpls = c.execute("SELECT t.*, "
                      "(SELECT COUNT(*) FROM ie_template_sections s WHERE s.template_id=t.id) AS n_sec, "
                      "(SELECT COUNT(*) FROM ie_template_items i JOIN ie_template_sections s "
                      " ON s.id=i.section_id WHERE s.template_id=t.id) AS n_items "
                      "FROM ie_templates t ORDER BY t.active DESC, t.name").fetchall()
    sel = t or (tmpls[0]["id"] if tmpls else 0)
    cur = c.execute("SELECT * FROM ie_templates WHERE id=?", (sel,)).fetchone()
    sections = []
    if cur:
        for s in c.execute("SELECT * FROM ie_template_sections WHERE template_id=? "
                           "ORDER BY sort_order, id", (sel,)):
            items = c.execute("SELECT * FROM ie_template_items WHERE section_id=? "
                              "ORDER BY sort_order, id", (s["id"],)).fetchall()
            sections.append({"sec": s, "items": items})
    history = c.execute("SELECT * FROM ie_template_revisions WHERE template_id=? "
                        "ORDER BY id DESC LIMIT 15", (sel,)).fetchall() if cur else []
    usage = c.execute("SELECT r.id, r.name, p.name AS product FROM ie_reports r "
                      "JOIN products p ON p.id=r.product_id WHERE r.template_id=? AND r.active=1",
                      (sel,)).fetchall() if cur else []
    c.close()
    return templates.TemplateResponse(request, "ie_templates.html", {"user": user, "tmpls": tmpls,
        "cur": cur, "sections": sections, "history": history, "usage": usage,
        "categories": db.CATEGORIES, "reviewers": db.IE_REVIEWERS, "priorities": db.IE_PRIORITIES})

@app.post("/ie/templates/add")
def ie_tmpl_add(request: Request, name: str = Form(...), reviewer: str = Form("DNV"),
                category: str = Form(""), notes: str = Form(""), source_url: str = Form(""),
                copy_from: str = Form(""), user=Depends(require_editor)):
    nm = name.strip()
    if not nm:
        return RedirectResponse("/ie/templates", status_code=303)
    c = db.conn()
    if c.execute("SELECT 1 FROM ie_templates WHERE name=?", (nm,)).fetchone():
        c.close()
        return RedirectResponse("/ie/templates?err=dup", status_code=303)
    c.execute("INSERT INTO ie_templates(name, reviewer, category, notes, source_url, "
              "created_by, created_at) VALUES(?,?,?,?,?,?,?)",
              (nm, reviewer, category if category in db.CATEGORIES else "", notes.strip(),
               source_url.strip(), user["email"], now()))
    tid = c.execute("SELECT id FROM ie_templates WHERE name=?", (nm,)).fetchone()["id"]
    if copy_from.strip().isdigit():
        for s in c.execute("SELECT * FROM ie_template_sections WHERE template_id=? "
                           "ORDER BY sort_order, id", (int(copy_from),)).fetchall():
            c.execute("INSERT INTO ie_template_sections(template_id, code, title, owner, sort_order) "
                      "VALUES(?,?,?,?,?)", (tid, s["code"], s["title"], s["owner"], s["sort_order"]))
            nsid = c.execute("SELECT last_insert_rowid() AS id").fetchone()["id"]
            c.execute("INSERT INTO ie_template_items(section_id, item_id, sub_section, review_item, "
                      "evidence, suggested_owner, priority, source, sort_order) "
                      "SELECT ?, item_id, sub_section, review_item, evidence, suggested_owner, "
                      "priority, source, sort_order FROM ie_template_items WHERE section_id=?",
                      (nsid, s["id"]))
    c.commit(); c.close()
    db.log(user["email"], "ie:template:add", nm)
    return RedirectResponse(f"/ie/templates?t={tid}", status_code=303)

@app.post("/ie/templates/{tid}/edit")
def ie_tmpl_edit(tid: int, request: Request, name: str = Form(...), reviewer: str = Form("DNV"),
                 category: str = Form(""), notes: str = Form(""), source_url: str = Form(""),
                 user=Depends(require_editor)):
    nm = name.strip()
    if not nm:
        return RedirectResponse(f"/ie/templates?t={tid}", status_code=303)
    c = db.conn()
    if c.execute("SELECT 1 FROM ie_templates WHERE name=? AND id<>?", (nm, tid)).fetchone():
        c.close()
        return RedirectResponse(f"/ie/templates?t={tid}&err=dup", status_code=303)
    db.ie_record_revision(c, tid, "scope", f"renamed/rescoped to '{nm}'", user["email"])
    c.execute("UPDATE ie_templates SET name=?, reviewer=?, category=?, notes=?, source_url=? "
              "WHERE id=?", (nm, reviewer, category if category in db.CATEGORIES else "",
                             notes.strip(), source_url.strip(), tid))
    c.commit(); c.close()
    db.log(user["email"], "ie:template:edit", nm)
    return RedirectResponse(f"/ie/templates?t={tid}", status_code=303)

@app.post("/ie/templates/{tid}/toggle")
def ie_tmpl_toggle(tid: int, request: Request, user=Depends(require_editor)):
    c = db.conn()
    db.ie_record_revision(c, tid, "retire/restore", "", user["email"])
    c.execute("UPDATE ie_templates SET active = 1 - active WHERE id=?", (tid,))
    c.commit(); c.close()
    return RedirectResponse(f"/ie/templates?t={tid}", status_code=303)

@app.post("/ie/templates/{tid}/section/add")
def ie_sec_add(tid: int, request: Request, title: str = Form(...), code: str = Form(""),
               owner: str = Form(""), user=Depends(require_editor)):
    ti = title.strip()
    if ti:
        c = db.conn()
        db.ie_record_revision(c, tid, "add section", ti, user["email"])
        nxt = c.execute("SELECT COALESCE(MAX(sort_order), 0) + 1 FROM ie_template_sections "
                        "WHERE template_id=?", (tid,)).fetchone()[0]
        c.execute("INSERT INTO ie_template_sections(template_id, code, title, owner, sort_order) "
                  "VALUES(?,?,?,?,?)", (tid, code.strip(), ti, owner.strip(), nxt))
        c.commit(); c.close()
        db.log(user["email"], "ie:section:add", ti)
    return RedirectResponse(f"/ie/templates?t={tid}", status_code=303)

@app.post("/ie/templates/section/{sid}/save")
def ie_sec_save(sid: int, request: Request, title: str = Form(...), code: str = Form(""),
                owner: str = Form(""), sort_order: int = Form(0), user=Depends(require_editor)):
    c = db.conn()
    row = c.execute("SELECT template_id, title FROM ie_template_sections WHERE id=?", (sid,)).fetchone()
    if not row or not title.strip():
        c.close()
        return RedirectResponse("/ie/templates", status_code=303)
    db.ie_record_revision(c, row["template_id"], "edit section", row["title"], user["email"])
    c.execute("UPDATE ie_template_sections SET title=?, code=?, owner=?, sort_order=? WHERE id=?",
              (title.strip(), code.strip(), owner.strip(), sort_order, sid))
    c.commit(); c.close()
    return RedirectResponse(f"/ie/templates?t={row['template_id']}", status_code=303)

@app.post("/ie/templates/section/{sid}/delete")
def ie_sec_delete(sid: int, request: Request, user=Depends(require_editor)):
    c = db.conn()
    row = c.execute("SELECT template_id, title FROM ie_template_sections WHERE id=?", (sid,)).fetchone()
    if not row:
        c.close()
        return RedirectResponse("/ie/templates", status_code=303)
    db.ie_record_revision(c, row["template_id"], "delete section", row["title"], user["email"])
    c.execute("DELETE FROM ie_template_sections WHERE id=?", (sid,))
    c.commit(); c.close()
    db.log(user["email"], "ie:section:delete", row["title"])
    return RedirectResponse(f"/ie/templates?t={row['template_id']}", status_code=303)

@app.post("/ie/templates/section/{sid}/item/add")
def ie_tmpl_item_add(sid: int, request: Request, review_item: str = Form(...),
                     item_id: str = Form(""), sub_section: str = Form(""),
                     evidence: str = Form(""), suggested_owner: str = Form(""),
                     priority: str = Form("Normal"), source: str = Form(""),
                     user=Depends(require_editor)):
    ri = review_item.strip()
    c = db.conn()
    row = c.execute("SELECT template_id FROM ie_template_sections WHERE id=?", (sid,)).fetchone()
    if not row or not ri:
        c.close()
        return RedirectResponse("/ie/templates", status_code=303)
    db.ie_record_revision(c, row["template_id"], "add item", ri, user["email"])
    nxt = c.execute("SELECT COALESCE(MAX(sort_order), 0) + 1 FROM ie_template_items "
                    "WHERE section_id=?", (sid,)).fetchone()[0]
    c.execute("INSERT INTO ie_template_items(section_id, item_id, sub_section, review_item, "
              "evidence, suggested_owner, priority, source, sort_order) VALUES(?,?,?,?,?,?,?,?,?)",
              (sid, item_id.strip(), sub_section.strip(), ri, evidence.strip(),
               suggested_owner.strip(),
               priority if priority in db.IE_PRIORITIES else "Normal", source.strip(), nxt))
    c.commit(); c.close()
    db.log(user["email"], "ie:template:item:add", ri[:60])
    return RedirectResponse(f"/ie/templates?t={row['template_id']}", status_code=303)

@app.post("/ie/templates/item/{iid}/save")
def ie_tmpl_item_save(iid: int, request: Request, review_item: str = Form(...),
                      item_id: str = Form(""), sub_section: str = Form(""),
                      evidence: str = Form(""), suggested_owner: str = Form(""),
                      priority: str = Form("Normal"), sort_order: int = Form(0),
                      user=Depends(require_editor)):
    c = db.conn()
    row = c.execute("SELECT i.review_item, s.template_id FROM ie_template_items i "
                    "JOIN ie_template_sections s ON s.id=i.section_id WHERE i.id=?", (iid,)).fetchone()
    if not row or not review_item.strip():
        c.close()
        return RedirectResponse("/ie/templates", status_code=303)
    db.ie_record_revision(c, row["template_id"], "edit item", row["review_item"], user["email"])
    c.execute("UPDATE ie_template_items SET review_item=?, item_id=?, sub_section=?, evidence=?, "
              "suggested_owner=?, priority=?, sort_order=? WHERE id=?",
              (review_item.strip(), item_id.strip(), sub_section.strip(), evidence.strip(),
               suggested_owner.strip(),
               priority if priority in db.IE_PRIORITIES else "Normal", sort_order, iid))
    c.commit(); c.close()
    return RedirectResponse(f"/ie/templates?t={row['template_id']}", status_code=303)

@app.post("/ie/templates/item/{iid}/delete")
def ie_tmpl_item_delete(iid: int, request: Request, user=Depends(require_editor)):
    c = db.conn()
    row = c.execute("SELECT i.review_item, s.template_id FROM ie_template_items i "
                    "JOIN ie_template_sections s ON s.id=i.section_id WHERE i.id=?", (iid,)).fetchone()
    if not row:
        c.close()
        return RedirectResponse("/ie/templates", status_code=303)
    db.ie_record_revision(c, row["template_id"], "delete item", row["review_item"], user["email"])
    c.execute("DELETE FROM ie_template_items WHERE id=?", (iid,))
    c.commit(); c.close()
    return RedirectResponse(f"/ie/templates?t={row['template_id']}", status_code=303)

@app.post("/ie/templates/undo/{rev_id}")
def ie_tmpl_undo(rev_id: int, request: Request, user=Depends(require_editor)):
    c = db.conn()
    res = db.ie_restore_revision(c, rev_id, user["email"])
    if not res:
        c.close()
        return RedirectResponse("/ie/templates?err=gone", status_code=303)
    tid, action = res
    c.commit(); c.close()
    db.log(user["email"], "ie:template:undo", f"t{tid}: reverted '{action}'")
    return RedirectResponse(f"/ie/templates?t={tid}&undone={action}", status_code=303)

@app.get("/ie/report/{rid}/bundle.zip")
def ie_bundle(rid: int, request: Request, user=Depends(require_user)):
    """The IE evidence pack: files foldered by section, with manifest and summary."""
    c = db.conn()
    r = c.execute("SELECT r.*, p.name AS product FROM ie_reports r JOIN products p ON p.id=r.product_id "
                  "WHERE r.id=?", (rid,)).fetchone()
    if not r:
        c.close()
        return RedirectResponse("/ie", status_code=303)
    groups, totals, atts, sec_atts = _ie_progress(c, rid)
    c.close()
    stamp = datetime.datetime.now().strftime("%Y%m%d")
    base = f"{_safe(r['reviewer'], 20)}_{_safe(r['product'], 40)}_IE_{stamp}"
    path = os.path.join(PKG_DIR, base + ".zip")
    man, summ = [], []
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        for gi, g in enumerate(groups, 1):
            folder = f"{gi:02d}_{_safe(g['sec']['title'], 50)}"
            # Section-scoped evidence sits at the top of its section folder,
            # above the per-item subfolders.
            for f in sec_atts.get(g["sec"]["id"], []):
                arc = f"{folder}/_section/{_safe(f['filename'], 90)}"
                if os.path.exists(f["stored_path"]):
                    z.write(f["stored_path"], arc)
                    man.append([g["sec"]["title"], "", "(section evidence)", arc])
                else:
                    man.append([g["sec"]["title"], "", "(section evidence)",
                                "FILE NOT FOUND ON DISK"])
            if g["sec"]["narrative"]:
                z.writestr(f"{folder}/_section/NARRATIVE.txt",
                           f"{g['sec']['title']}\n{'=' * len(g['sec']['title'])}\n\n"
                           f"{g['sec']['narrative']}\n")
                man.append([g["sec"]["title"], "", "(section narrative)",
                            f"{folder}/_section/NARRATIVE.txt"])
            for it in g["items"]:
                files = atts.get(it["id"], [])
                summ.append([g["sec"]["title"], it["item_id"], it["sub_section"],
                             it["review_item"], it["evidence"], it["priority"], it["status"],
                             it["owner"], it["due_date"], len(files),
                             "INCLUDED" if files else "NO EVIDENCE", it["gap"]])
                for f in files:
                    arc = f"{folder}/{_safe(it['item_id'] or it['sub_section'], 40)}/{_safe(f['filename'], 90)}"
                    if os.path.exists(f["stored_path"]):
                        z.write(f["stored_path"], arc)
                        man.append([g["sec"]["title"], it["item_id"], it["review_item"], arc])
                    else:
                        man.append([g["sec"]["title"], it["item_id"], it["review_item"],
                                    "FILE NOT FOUND ON DISK"])

        def csv_bytes(header, rows):
            buf = io.StringIO(); w = csv.writer(buf); w.writerow(header)
            for row in rows:
                w.writerow(row)
            return buf.getvalue()

        narr = [g for g in groups if g["sec"]["narrative"]]
        if narr:
            z.writestr("IE_NARRATIVE.txt", "\n\n".join(
                f"{g['sec']['title']}\n{'=' * len(g['sec']['title'])}\n"
                f"Owner: {g['sec']['owner'] or '-'}\n\n{g['sec']['narrative']}"
                for g in narr) + "\n")

        z.writestr("MANIFEST.csv", csv_bytes(
            ["Section", "Item ID", "Review item", "Path in zip"], man))
        z.writestr("IE_SUMMARY.csv", csv_bytes(
            ["Section", "Item ID", "Sub-section", "Preparation item / review question",
             "Required evidence", "Priority", "Status", "Owner", "Due", "Files", "Result",
             "Gap / action"], summ))
        L = ["IE TECHNOLOGY REVIEW - EVIDENCE SUMMARY", "=" * 62,
             f"Product      : {r['product']}",
             f"Report       : {r['name']}",
             f"Reviewer     : {r['reviewer']}",
             f"Status       : {r['status']}",
             f"Kickoff      : {r['kickoff_date'] or '-'}    Target: {r['target_date'] or '-'}",
             f"Built        : {now()} by {user['email']}", ""]
        if r["shared_with"]:
            L += [f"Shared with  : {r['shared_with']}", ""]
        L += [f"{totals['done']}/{totals['n']} items accepted ({totals['pct']}%), "
              f"{totals['wip']} in progress, {totals['blocked']} blocked, "
              f"{totals['crit_open']} Critical still open.",
              f"{totals['files']} evidence file(s) attached.", "",
              "BY SECTION", "-" * 62,
              f"{'Section':<44}{'Items':>6}{'Done':>6}{'Files':>6}"]
        for g in groups:
            L.append(f"{g['sec']['title'][:43]:<44}{g['n']:>6}{g['done']:>6}{g['files']:>6}")
        L += [f"{'TOTAL':<44}{totals['n']:>6}{totals['done']:>6}{totals['files']:>6}", "",
              "DETAIL", "-" * 62]
        for gi, g in enumerate(groups, 1):
            L.append(f"{gi:02d}. {g['sec']['title']}  ({g['done']}/{g['n']}, owner: {g['sec']['owner'] or '-'})")
            for it in g["items"]:
                files = atts.get(it["id"], [])
                mark = "[x]" if files else "[ ]"
                L.append(f"    {mark} {it['item_id']:<6} ({it['priority']}, {it['status']}) {it['review_item']}")
                if it["evidence"]:
                    L.append(f"           evidence: {it['evidence']}")
                for f in files:
                    L.append(f"           - {f['filename']}")
                if it["gap"]:
                    L.append(f"           GAP: {it['gap']}")
            L.append("")
        z.writestr("IE_SUMMARY.txt", "\n".join(L))
    db.log(user["email"], "ie:bundle", f"{r['name']} - {totals['files']} files")
    return FileResponse(path, filename=base + ".zip", media_type="application/zip")

# ---------------- 13) AVL acceptance: the pursuit of one product at one TPO ----------------
# Everything else is organised by function; this is organised by the thing being
# pursued, and only reads what the other pages already store.

def _pursuit_readiness(c, avl_id, product_id):
    """Dataroom readiness for one product x TPO, by obligation."""
    rows = c.execute("SELECT * FROM checklist_items WHERE product_id=? AND avl_id=?",
                     (product_id, avl_id)).fetchall()
    ids = [r["id"] for r in rows]
    with_file = set()
    if ids:
        qs = ",".join("?" * len(ids))
        with_file = {r[0] for r in c.execute(
            f"SELECT DISTINCT ref_id FROM attachments WHERE kind='checklist' AND ref_id IN ({qs})", ids)}
    req = [r for r in rows if r["obligation"] == "Required"]
    done = [r for r in req if r["status"] in db.CHECK_DONE]
    return {"n": len(rows), "req": len(req), "done": len(done),
            "accepted": sum(1 for r in req if r["status"] == "Accepted"),
            "blocked": sum(1 for r in rows if r["status"] == "Blocked"),
            "nofile_req": sum(1 for r in req if r["id"] not in with_file),
            "files": len(with_file),
            "pct": round(100 * len(done) / len(req)) if req else None,
            "ready": bool(req) and len(done) == len(req)}

def _pursuit_row(c, avl_id, product_id):
    """One line of the portfolio: status, readiness, package, actions, IE, last activity."""
    li = c.execute("SELECT * FROM listings WHERE avl_id=? AND product_id=?",
                   (avl_id, product_id)).fetchone()
    rd = _pursuit_readiness(c, avl_id, product_id)
    pkg = c.execute("SELECT rev_code, rev_date, n_gaps FROM packages WHERE avl_id=? AND product_id=? "
                    "ORDER BY revision DESC LIMIT 1", (avl_id, product_id)).fetchone()
    today = datetime.date.today().isoformat()
    acts = c.execute("SELECT COUNT(*) n, SUM(CASE WHEN COALESCE(due_date,'')<>'' AND due_date<? "
                     "THEN 1 ELSE 0 END) od FROM actions WHERE avl_id=? AND product_id=? "
                     "AND status='Open'", (today, avl_id, product_id)).fetchone()
    ie = c.execute("SELECT r.id, r.name, r.status, r.reviewer, "
                   "(SELECT COUNT(*) FROM ie_report_items i WHERE i.report_id=r.id) n, "
                   "(SELECT COUNT(*) FROM ie_report_items i WHERE i.report_id=r.id "
                   " AND i.status IN ('Accepted','N/A')) done "
                   "FROM ie_reports r WHERE r.product_id=? AND r.active=1 "
                   "ORDER BY r.id DESC LIMIT 1", (product_id,)).fetchone()
    last = c.execute("SELECT MAX(ts) t FROM (SELECT ts FROM audit WHERE avl_id=? AND product_id=? "
                     "UNION ALL SELECT ts FROM status_history WHERE avl_id=? AND product_id=?)",
                     (avl_id, product_id, avl_id, product_id)).fetchone()["t"]
    nxt = c.execute("SELECT * FROM commitments WHERE avl_id=? AND product_id=? "
                    "AND status='Planned' ORDER BY due_date LIMIT 1",
                    (avl_id, product_id)).fetchone()
    return {"listing": li, "rd": rd, "pkg": pkg, "ie": ie,
            "next_commit": nxt,
            "next_state": db.commitment_state(nxt, today) if nxt else None,
            "open_actions": acts["n"] or 0, "overdue_actions": acts["od"] or 0,
            "last_activity": last,
            "overdue_target": bool(li and li["target_date"] and li["target_date"] < today
                                   and li["status"] not in db.LISTED_STATUSES)}

@app.get("/acceptance", response_class=HTMLResponse)
def acceptance(request: Request, avl: int = 0, view: str = "active", user=Depends(require_user)):
    """Portfolio of every product x TPO being pursued."""
    c = db.conn()
    avls = c.execute("SELECT id, name FROM avls WHERE active=1 ORDER BY name").fetchall()
    q = ("SELECT l.*, a.name AS avl_name, p.name AS product, p.category, p.lifecycle "
         "FROM listings l JOIN avls a ON a.id=l.avl_id AND a.active=1 "
         "JOIN products p ON p.id=l.product_id AND p.active=1 WHERE 1=1 ")
    args = []
    if avl:
        q += "AND l.avl_id=? "; args.append(avl)
    if view == "active":
        q += "AND l.status NOT IN ('N/A','No Interest') "
    elif view == "listed":
        q += "AND l.status IN ('Listed','Listed, Conditional') "
    elif view == "conditional":
        q += "AND l.status='Listed, Conditional' "
    elif view == "pursuing":
        q += "AND l.status NOT IN ('N/A','No Interest','Listed') "
    rows = []
    for r in c.execute(q + "ORDER BY a.name, p.category, p.name", args):
        d = _pursuit_row(c, r["avl_id"], r["product_id"])
        rows.append({"l": r, **d})
    people = c.execute("SELECT id, name FROM people WHERE active=1 ORDER BY name").fetchall()
    c.close()
    tot = {
        "n": len(rows),
        "listed": sum(1 for r in rows if r["l"]["status"] in db.LISTED_STATUSES),
        "conditional": sum(1 for r in rows if r["l"]["status"] == "Listed, Conditional"),
        "ready": sum(1 for r in rows if r["rd"]["ready"]),
        "blocked": sum(1 for r in rows if r["rd"]["blocked"]),
        "overdue": sum(1 for r in rows if r["overdue_target"]),
    }
    return templates.TemplateResponse(request, "acceptance.html", {"user": user, "rows": rows,
        "avls": avls, "sel_a": avl, "view": view, "tot": tot, "people": people})

@app.get("/pursuit/{avl_id}/{product_id}", response_class=HTMLResponse)
def pursuit(avl_id: int, product_id: int, request: Request, user=Depends(require_user)):
    c = db.conn()
    a = c.execute("SELECT * FROM avls WHERE id=?", (avl_id,)).fetchone()
    p = c.execute("SELECT * FROM products WHERE id=?", (product_id,)).fetchone()
    if not a or not p:
        c.close()
        return RedirectResponse("/acceptance", status_code=303)
    # Create the listing lazily so a pursuit page always exists for a real pair.
    if not c.execute("SELECT 1 FROM listings WHERE avl_id=? AND product_id=?",
                     (avl_id, product_id)).fetchone():
        c.execute("INSERT INTO listings(product_id, avl_id, status, updated_by, updated_at) "
                  "VALUES(?,?,?,?,?)", (product_id, avl_id, "No Info", user["email"], now()))
        c.commit()
    d = _pursuit_row(c, avl_id, product_id)
    today = datetime.date.today().isoformat()
    gaps = c.execute("SELECT * FROM checklist_items WHERE product_id=? AND avl_id=? "
                     "AND obligation='Required' AND status NOT IN ('Complete','Submitted','Accepted') "
                     "ORDER BY sort_order, id LIMIT 25", (product_id, avl_id)).fetchall()
    acts = c.execute("SELECT ac.*, pe.name AS owner_name, ci.workstream FROM actions ac "
                     "LEFT JOIN people pe ON pe.id=ac.owner_person_id "
                     "LEFT JOIN checklist_items ci ON ci.id=ac.checklist_item_id "
                     "WHERE ac.avl_id=? AND ac.product_id=? ORDER BY ac.status, "
                     "CASE WHEN COALESCE(ac.due_date,'')='' THEN 1 ELSE 0 END, ac.due_date",
                     (avl_id, product_id)).fetchall()
    calls_ = c.execute("SELECT * FROM calls WHERE avl_id=? ORDER BY call_date DESC LIMIT 5",
                       (avl_id,)).fetchall()
    contacts_ = c.execute("SELECT * FROM contacts WHERE avl_id=? AND active=1 "
                          "ORDER BY is_primary DESC, name", (avl_id,)).fetchall()
    packages_ = c.execute("SELECT * FROM packages WHERE avl_id=? AND product_id=? "
                          "ORDER BY revision DESC LIMIT 5", (avl_id, product_id)).fetchall()
    commits = _commitment_rows(c, "WHERE cm.avl_id=? AND cm.product_id=? ",
                               (avl_id, product_id), limit=20)
    feed = _activity(c, avl=avl_id, product=product_id, limit=20)
    people = c.execute("SELECT id, name, org FROM people WHERE active=1 ORDER BY name").fetchall()
    c.close()
    return templates.TemplateResponse(request, "pursuit.html", {"user": user, "a": a, "p": p,
        "d": d, "l": d["listing"], "gaps": gaps, "acts": acts, "calls": calls_,
        "contacts": contacts_, "packages": packages_, "feed": feed, "people": people,
        "commits": commits, "kinds": db.COMMITMENT_KINDS,
        "today": today, "statuses": db.STATUSES, "priorities": db.PURSUIT_PRIORITIES,
        "types": db.ACTIVITY_TYPES})

@app.post("/pursuit/{avl_id}/{product_id}/save")
def pursuit_save(avl_id: int, product_id: int, request: Request,
                 owner_person_id: str = Form(""), owner_other: str = Form(""),
                 target_date: str = Form(""), submitted_at: str = Form(""),
                 condition: str = Form(""), next_milestone: str = Form(""),
                 risk: str = Form(""), priority: str = Form("Normal"),
                 user=Depends(require_editor)):
    if priority not in db.PURSUIT_PRIORITIES:
        priority = "Normal"
    c = db.conn()
    oid = _as_id(owner_person_id)
    oname = ""
    if oid:
        row = c.execute("SELECT name FROM people WHERE id=? AND active=1", (oid,)).fetchone()
        oname, oid = (row["name"], oid) if row else ("", None)
    prev = c.execute("SELECT owner FROM listings WHERE avl_id=? AND product_id=?",
                     (avl_id, product_id)).fetchone()
    owner_txt = oname or owner_other.strip() or (prev["owner"] if prev and not oid else "")
    c.execute("UPDATE listings SET owner=?, owner_person_id=?, target_date=?, submitted_at=?, "
              "condition=?, next_milestone=?, risk=?, priority=? WHERE avl_id=? AND product_id=?",
              (owner_txt, oid, target_date, submitted_at, condition.strip(),
               next_milestone.strip(), risk.strip(), priority, avl_id, product_id))
    c.commit(); c.close()
    db.log(user["email"], "pursuit:save", next_milestone.strip()[:60] or "details updated",
           avl_id=avl_id, product_id=product_id, entity="listing")
    return RedirectResponse(f"/pursuit/{avl_id}/{product_id}", status_code=303)

# ---------------- 14) schedule: dated commitments ----------------
def _commitment_rows(c, where="", args=(), limit=300):
    """Commitments with the pursuit's readiness attached, so a date can be read
    against what is actually outstanding rather than on its own."""
    rows = c.execute(
        "SELECT cm.*, a.name AS avl_name, p.name AS product, p.category, "
        "pe.name AS owner_name, l.status AS listing_status "
        "FROM commitments cm JOIN avls a ON a.id=cm.avl_id "
        "JOIN products p ON p.id=cm.product_id "
        "LEFT JOIN people pe ON pe.id=cm.owner_person_id "
        "LEFT JOIN listings l ON l.avl_id=cm.avl_id AND l.product_id=cm.product_id "
        f"{where} ORDER BY CASE cm.status WHEN 'Planned' THEN 0 ELSE 1 END, "
        "cm.due_date, a.name LIMIT ?", list(args) + [limit]).fetchall()
    today = datetime.date.today().isoformat()
    out = []
    for r in rows:
        rd = _pursuit_readiness(c, r["avl_id"], r["product_id"])
        n_open = c.execute("SELECT COUNT(*) FROM actions WHERE avl_id=? AND product_id=? "
                           "AND status='Open'", (r["avl_id"], r["product_id"])).fetchone()[0]
        out.append({"c": r, "rd": rd, "open_actions": n_open,
                    **db.commitment_state(r, today)})
    return out

@app.get("/schedule", response_class=HTMLResponse)
def schedule(request: Request, show: str = "open", avl: int = 0, owner: int = 0,
             user=Depends(require_user)):
    c = db.conn()
    where, args = "WHERE 1=1 ", []
    if show == "open":
        where += "AND cm.status='Planned' "
    elif show in ("Met", "Missed", "Cancelled"):
        where += "AND cm.status=? "; args.append(show)
    if avl:
        where += "AND cm.avl_id=? "; args.append(avl)
    if owner:
        where += "AND cm.owner_person_id=? "; args.append(owner)
    rows = _commitment_rows(c, where, args)
    avls = c.execute("SELECT id, name FROM avls WHERE active=1 ORDER BY name").fetchall()
    products = c.execute("SELECT id, name, category FROM products WHERE active=1 "
                         "ORDER BY category, name").fetchall()
    people = c.execute("SELECT id, name FROM people WHERE active=1 ORDER BY name").fetchall()
    c.close()
    tot = {"open": sum(1 for r in rows if r["c"]["status"] == "Planned"),
           "overdue": sum(1 for r in rows if r["state"] == "overdue"),
           "soon": sum(1 for r in rows if r["state"] == "due-soon"),
           "notready": sum(1 for r in rows if r["state"] in ("overdue", "due-soon")
                           and not r["rd"]["ready"])}
    return templates.TemplateResponse(request, "schedule.html", {"user": user, "rows": rows,
        "avls": avls, "products": products, "people": people, "show": show, "sel_a": avl,
        "sel_o": owner, "tot": tot, "kinds": db.COMMITMENT_KINDS,
        "statuses": db.COMMITMENT_STATUSES, "today": datetime.date.today().isoformat(),
        "risk_days": db.AT_RISK_DAYS})

@app.post("/schedule/add")
def commitment_add(request: Request, avl_id: int = Form(...), product_id: int = Form(...),
                   due_date: str = Form(...), kind: str = Form("Dataroom submission"),
                   label: str = Form(""), owner_person_id: str = Form(""),
                   owner_other: str = Form(""), notes: str = Form(""),
                   next_url: str = Form(""), user=Depends(require_editor)):
    dest = next_url if next_url.startswith("/") else "/schedule"
    if not due_date.strip():
        return RedirectResponse(dest + "?err=date", status_code=303)
    if kind not in db.COMMITMENT_KINDS:
        kind = "Other"
    c = db.conn()
    oid = _as_id(owner_person_id)
    oname = ""
    if oid:
        row = c.execute("SELECT name FROM people WHERE id=? AND active=1", (oid,)).fetchone()
        oname, oid = (row["name"], oid) if row else ("", None)
    c.execute("INSERT INTO commitments(avl_id, product_id, kind, label, due_date, owner, "
              "owner_person_id, notes, created_by, created_at) VALUES(?,?,?,?,?,?,?,?,?,?)",
              (avl_id, product_id, kind, label.strip(), due_date,
               oname or owner_other.strip(), oid, notes.strip(), user["email"], now()))
    names = c.execute("SELECT a.name AS a, p.name AS p FROM avls a, products p "
                      "WHERE a.id=? AND p.id=?", (avl_id, product_id)).fetchone()
    c.commit(); c.close()
    db.log(user["email"], "commitment:add", f"{kind} by {due_date}",
           avl_id=avl_id, product_id=product_id, entity="commitment")
    return RedirectResponse(dest, status_code=303)

@app.post("/schedule/{cid}/save")
def commitment_save(cid: int, request: Request, due_date: str = Form(...),
                    kind: str = Form("Other"), label: str = Form(""),
                    owner_person_id: str = Form(""), owner_other: str = Form(""),
                    notes: str = Form(""), next_url: str = Form(""),
                    user=Depends(require_editor)):
    dest = next_url if next_url.startswith("/") else "/schedule"
    c = db.conn()
    row = c.execute("SELECT * FROM commitments WHERE id=?", (cid,)).fetchone()
    if not row or not due_date.strip():
        c.close()
        return RedirectResponse(dest, status_code=303)
    oid = _as_id(owner_person_id)
    oname = ""
    if oid:
        p = c.execute("SELECT name FROM people WHERE id=? AND active=1", (oid,)).fetchone()
        oname, oid = (p["name"], oid) if p else ("", None)
    owner_txt = oname or owner_other.strip() or (row["owner"] if not oid else "")
    c.execute("UPDATE commitments SET due_date=?, kind=?, label=?, owner=?, owner_person_id=?, "
              "notes=? WHERE id=?",
              (due_date, kind if kind in db.COMMITMENT_KINDS else "Other", label.strip(),
               owner_txt, oid, notes.strip(), cid))
    c.commit(); c.close()
    db.log(user["email"], "commitment:save", f"{kind} -> {due_date}",
           avl_id=row["avl_id"], product_id=row["product_id"], entity="commitment")
    return RedirectResponse(dest, status_code=303)

@app.post("/schedule/{cid}/status")
def commitment_status(cid: int, request: Request, status: str = Form(...),
                      next_url: str = Form(""), user=Depends(require_editor)):
    dest = next_url if next_url.startswith("/") else "/schedule"
    if status not in db.COMMITMENT_STATUSES:
        return RedirectResponse(dest, status_code=303)
    c = db.conn()
    row = c.execute("SELECT * FROM commitments WHERE id=?", (cid,)).fetchone()
    if not row:
        c.close()
        return RedirectResponse(dest, status_code=303)
    met = now()[:10] if status == "Met" else ""
    c.execute("UPDATE commitments SET status=?, met_at=? WHERE id=?", (status, met, cid))
    # Meeting a dataroom submission is the event the pursuit records as submitted.
    if status == "Met" and row["kind"] == "Dataroom submission":
        c.execute("UPDATE listings SET submitted_at=? WHERE avl_id=? AND product_id=? "
                  "AND COALESCE(submitted_at,'')=''", (met, row["avl_id"], row["product_id"]))
    c.commit(); c.close()
    db.log(user["email"], "commitment:status", f"{row['kind']} ({row['due_date']}) -> {status}",
           avl_id=row["avl_id"], product_id=row["product_id"], entity="commitment")
    return RedirectResponse(dest, status_code=303)

@app.post("/schedule/{cid}/delete")
def commitment_delete(cid: int, request: Request, next_url: str = Form(""),
                      user=Depends(require_editor)):
    dest = next_url if next_url.startswith("/") else "/schedule"
    c = db.conn()
    row = c.execute("SELECT * FROM commitments WHERE id=?", (cid,)).fetchone()
    c.execute("DELETE FROM commitments WHERE id=?", (cid,))
    c.commit(); c.close()
    if row:
        db.log(user["email"], "commitment:delete", f"{row['kind']} {row['due_date']}",
               avl_id=row["avl_id"], product_id=row["product_id"], entity="commitment")
    return RedirectResponse(dest, status_code=303)

@app.post("/schedule/{cid}/cascade")
def commitment_cascade(cid: int, request: Request, lead_days: int = Form(7),
                       overwrite: int = Form(0), next_url: str = Form(""),
                       user=Depends(require_editor)):
    """Work backwards: give the outstanding required requirements a due date.

    Defaults to filling only blanks, because someone who has already dated a
    requirement knows something this calculation does not.
    """
    dest = next_url if next_url.startswith("/") else "/schedule"
    c = db.conn()
    row = c.execute("SELECT * FROM commitments WHERE id=?", (cid,)).fetchone()
    if not row:
        c.close()
        return RedirectResponse(dest, status_code=303)
    try:
        target = (datetime.date.fromisoformat(row["due_date"])
                  - datetime.timedelta(days=max(0, lead_days))).isoformat()
    except ValueError:
        c.close()
        return RedirectResponse(dest + "?err=date", status_code=303)
    q = ("UPDATE checklist_items SET due_date=?, updated_by=?, updated_at=? "
         "WHERE product_id=? AND avl_id=? AND obligation='Required' "
         "AND status NOT IN ('Complete','Submitted','Accepted') ")
    if not overwrite:
        q += "AND COALESCE(due_date,'')='' "
    c.execute(q, (target, user["email"], now(), row["product_id"], row["avl_id"]))
    n = c.total_changes
    c.commit(); c.close()
    db.log(user["email"], "commitment:cascade",
           f"{n} requirement(s) dated {target} for the {row['due_date']} commitment",
           avl_id=row["avl_id"], product_id=row["product_id"], entity="checklist")
    return RedirectResponse(dest + ("&" if "?" in dest else "?") + f"ok=dated{n}", status_code=303)
